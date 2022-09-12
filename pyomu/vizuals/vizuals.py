
import pandas as pd
from pandas._libs.lib import is_integer
import numpy as np
import geopandas as gpd
import h3
import matplotlib.pyplot as plt
import seaborn as sns

import mapclassify
import contextily as ctx

from tqdm.notebook import tqdm
tqdm.pandas()

import time

import unidecode
import itertools
from pathlib import Path

from sklearn.cluster import DBSCAN
from sklearn import metrics
from sklearn.metrics import silhouette_samples, silhouette_score
from sklearn.preprocessing import StandardScaler, RobustScaler 
from sklearn.decomposition import PCA

import osmnx as ox
import networkx as nx
import pandana

from pandana.loaders import osm as osm_pandana

import datetime as datetime
from datetime import date
from datetime import timedelta

pd.set_option('display.max_columns', None)

import googlemaps

import pytz
from tzwhere import tzwhere 
from pytz import timezone

import warnings
warnings.filterwarnings("ignore")
warnings.filterwarnings("ignore", category=DeprecationWarning) 


from shapely.geometry import Point, Polygon

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor # ColorFormat, 
from PIL import Image, ImageDraw

import os
import glob

import matplotlib as mpl
from mpl_toolkits.axes_grid1 import make_axes_locatable
import matplotlib.pyplot as plt
import matplotlib.ticker as tick
import IPython
from IPython.display import set_matplotlib_formats


from PIL import Image # pip install Pillow
from PIL import ImageOps

import seaborn as sns
from pandas import DataFrame

from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas
from matplotlib.figure import Figure
from matplotlib import cm
from matplotlib.colors import ListedColormap, LinearSegmentedColormap

from mpl_toolkits.axes_grid1 import make_axes_locatable


from pyomu.utils import utils


def crop_imagen(filePath, reduce=1, altura_max=0, ancho_max=0, save=True, crop_left = 0, crop_top = 0, crop_right = 0, crop_bottom = 0):
    
    # Trim all png images with white background in a folder
    # Usage "python PNGWhiteTrim.py ../someFolder padding"

    image=Image.open(filePath)
    image.load()
    imageSize = image.size #tuple
    
    ## QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    # remove alpha channel
    invert_im = image.convert("RGB")
    # invert image (so that white is 0)
    invert_im = ImageOps.invert(invert_im)
    imageBox = invert_im.getbbox()
    cropped=image.crop(imageBox)
    ## FIN DE QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    
    #REDUCE TAMAÑO
    _size=[]
    # calculates percentage to reduce image by maintaining proportion
    if altura_max>0: _size.append((altura_max/(cropped.height/38)))
    if ancho_max>0: _size.append((ancho_max/(cropped.width/38)))
    if len(_size) > 0: reduce = min(_size)
    
    if reduce < 1:
        basewidth = int(cropped.width * reduce)
        wpercent = (basewidth/float(cropped.size[0]))
        hsize = int((float(cropped.size[1])*float(wpercent)))
        # cropped.resize actually does the resizing
        cropped = cropped.resize((basewidth,hsize), Image.ANTIALIAS)
    
    if crop_left + crop_top + crop_right + crop_bottom > 0:
        width, height = cropped.size 
        crop_right = width - crop_right 
        crop_bottom = height - crop_bottom            
        cropped=cropped.crop((crop_left, crop_top, crop_right, crop_bottom))

    # save the image as cropped
    if save:
        filePath = filePath[0: filePath.find('.')]+'_cropped'+filePath[filePath.find('.'):len(filePath)]
        cropped.save(filePath)
        return filePath
    else:
        return cropped

def pptx_addtitle(prs, slide='', title='', top=0, left=0, width=10, height=1, new=True, fontsize=24, fontcolor='blue', bold=True):

    blank_slide_layout = prs.slide_layouts[6] # Using layout 6 (blank layout)
    # if new create blank slide
    if new:
        slide = prs.slides.add_slide(blank_slide_layout)

    # Set the slides background colour
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(212, 218, 220) # RGBColor(212, 218, 220) is the color of water on the contextily tiles

    # translates from cm to inches
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    # adds a text box onto the slide object
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Gill Sans'
    p.font.color.rgb = RGBColor(64,64,64) # (105,105,105) CSS Dim Grey
    if bold is True:
        p.font.bold = True
    p.font.size = Pt(fontsize)
    p.alignment = PP_ALIGN.CENTER
    
    #p.font.color = fontcolor
    # many more parameters available

    return slide

def pptx_addtext(prs, slide='', text='', top= 0, left=0, width=10, height=1):
    blank_slide_layout = prs.slide_layouts[6]
    
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.title = text
    # p.font.bold = True
    
    p.alignment = PP_ALIGN.RIGHT
    
    return slide

def pptx_addpic(prs, slide, img_path,  left=0, top=0, width=0, altura_max=0, ancho_max=0, crop_left = 0, crop_top = 0, crop_right = 0, crop_bottom = 0):
    # for adding all maps and graphs
    # altura_max and ancho_max in cm
    blank_slide_layout = prs.slide_layouts[6]

    img_path = str(img_path)

    if os.path.exists(img_path):
        # crop_imagen crops the image
        # NB commented out 20200514
        img_path = crop_imagen(img_path, reduce=1, altura_max=altura_max, ancho_max=ancho_max, save=True, crop_left=crop_left, crop_top=crop_top, crop_right=crop_right, crop_bottom=crop_bottom)
        
        # control position
        left = Inches(left)
        top = Inches(top)
        width  = Inches(width)
        # add to the slide
        if width!=0:
            slide_return = slide.shapes.add_picture(img_path, left, top, width) 
        else:
            slide_return = slide.shapes.add_picture(img_path, left, top) 
        
        os.remove(img_path)
        
        return slide_return

    
def indicators_all_day(od_matrix_sample, current_path = Path(), city='', prs='', title_ppt=''):
    '''
    Calcula indicadores para una muestra de viajes en automovil que se calcula para el día completo. Puede incluir varios días en el mismo DataFrame (ej. día de semana, sábado y domingo)
    od_matrix_sample = DataFrame con la matriz resultado de los viajes en automovil para las distintas horas de o los días.
    current_path = Directorio de trabajo. Por defecto Path()
    Salida. Dataframe de indicadores. Se guardan los resultados la carpeta de Resultados
    
    '''


    
    utils.create_result_dir(current_path)
    
    if 'driving_duration_in_traffic' in od_matrix_sample.columns:
    
        if len(od_matrix_sample[od_matrix_sample.driving_duration_in_traffic.notna()]) > 0:

            weekDaysMapping = ("Lunes", 
                               "Martes",
                               "Miércoles", 
                               "Jueves",
                               "Viernes", 
                               "Sábado",
                               "Domingo")
            monthsMapping = ("", 
                             "enero", 
                             "febrero", 
                             "marzo", 
                             "abril", 
                             "mayo", 
                             "junio", 
                             "julio", 
                             "agosto", 
                             "septiembre", 
                             "octubre", 
                             "noviembre", 
                             "diciembre")

            od_matrix_sample['kmh'] = (od_matrix_sample['driving_distance'] / (od_matrix_sample['driving_duration_in_traffic'] / 60)).round(2)

            trip_time = od_matrix_sample.groupby('trip_datetime', as_index=False).agg({'driving_duration_in_traffic':'mean', 'kmh': 'mean'})
            trip_time['date'] = trip_time.trip_datetime.dt.date
            trip_time['hour'] = trip_time['trip_datetime'].dt.hour.astype(str).str.zfill(2)

            indicators_all = pd.DataFrame([])
            for i in trip_time.date.unique():

                trip_time_tmp = trip_time[trip_time.date == i]
                _ = ''
                for x in range(0, len(f'{weekDaysMapping[i.weekday()]} {i.day} de {monthsMapping[i.month]}')):
                    _ += '-'

                indicadores = pd.DataFrame([[pd.to_datetime(i).date(), 
                                          f'{weekDaysMapping[i.weekday()]} {i.day} de {monthsMapping[i.month]}',
                                          trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmax()].hour,              
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmax()].driving_duration_in_traffic, 2),
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmax()].kmh, 2),
                                          trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmin()].hour,
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmin()].driving_duration_in_traffic, 2),
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmin()].kmh, 2),    
                                          round(trip_time_tmp.loc[trip_time_tmp.kmh.idxmax()].kmh / trip_time_tmp.loc[trip_time_tmp.kmh.idxmin()].kmh, 2),
                                          round(trip_time_tmp.driving_duration_in_traffic.mean(), 2),
                                          round(trip_time_tmp.kmh.mean(), 2)
                                          ]], 

                                         columns=['Date', 
                                                  'Detalle día',
                                                 'Hora Punta',                      
                                                 'Tiempo de viaje en hora punta (min)',
                                                 'Velocidad de viaje en hora punta (kmh)',
                                                 'Hora Valle', 
                                                 'Tiempo de viaje en hora valle (min)',
                                                 'Velocidad de viaje en hora valle (kmh)',
                                                 'Índice de congestión',
                                                 'Tiempo promedio de los viajes (min)', 
                                                 'Velocidad promedio de los viajes (kmh)', ])

                indicators_all = pd.concat([indicadores, indicators_all], ignore_index=True)

            df = indicators_all.set_index('Detalle día').T.reset_index().rename(columns={'index':'Detalle día'}).copy()

            fig = Figure(figsize=(6,2.5), dpi=300)
            canvas = FigureCanvas(fig)
            ax = fig.add_subplot(111)

            #hide the axes
            fig.patch.set_visible(False)
            ax.axis('off')

            colColours = ["#d4dadc" for i in range(0, len(df.columns))]
            cellColours = [colColours for i in range(0, len(df))]

            table = ax.table(cellText=df.values, 
                             colLabels=df.columns, 
                             loc='center', 
                             colLoc='center', 
                             cellLoc='center', 
                             cellColours=cellColours, 
                             colColours=colColours, 
                             colWidths=[.40,.17,.17,.17])

            table.set_fontsize(12)

            #display table
            fig.tight_layout()
            fig.savefig(current_path / 'Resultados_png' / f'{city}_indicadores_dia_completo.png', dpi=300)
            fig.savefig(current_path / 'Resultados_pdf' / f'{city}_indicadores_dia_completo.pdf', dpi=300)

            indicadores.to_csv(current_path / 'Resultados_files' / 'indicadores_dia_completo.csv', index=False)

            with sns.axes_style('darkgrid', {"axes.facecolor": "#d4dadc", 'figure.facecolor': "#d4dadc"}):

                # Tiempos promedios
                fig = Figure(figsize=(6,3), dpi=100)
                canvas = FigureCanvas(fig)
                ax = fig.add_subplot(111)


                for i in trip_time.date.unique():
                    var = f'{weekDaysMapping[i.weekday()]} {i.day}/{i.month}/{i.year}'
                    current_day = trip_time.loc[trip_time.date==i, ['hour', 'driving_duration_in_traffic']].reset_index(drop=True).rename(columns={'driving_duration_in_traffic':var})        
                    current_day.plot(ax=ax, legend=True, lw=.6)


                ax.set_title(f'Tiempos promedio de viaje', fontsize=8)
                ax.set_xlabel('Hora', fontsize=8)
                ax.set_ylabel('Tiempo promedio\n(minutos)', fontsize=8)
                ax.set_xticks(list(range(0, 24)))
                ax.tick_params(labelsize=6);

                box = ax.get_position()
                ax.set_position([box.x0, box.y0, box.width * 0.87, box.height])
                ax.legend(loc='center left', bbox_to_anchor=(1, 0.5), fontsize=6)
                fig.savefig(current_path / 'Resultados_png'  / f'{city}_tiempos_dia_completo.png', dpi=300)
                fig.savefig(current_path / 'Resultados_pdf'  / f'{city}_tiempos_dia_completo.pdf', dpi=300)
                if prs == '': display(fig)


                # Velocidades promedio
                fig = Figure(figsize=(6,3), dpi=100)
                canvas = FigureCanvas(fig)
                ax = fig.add_subplot(111)


                for i in trip_time.date.unique():
                    var = f'{weekDaysMapping[i.weekday()]} {i.day}/{i.month}/{i.year}'
                    current_day = trip_time.loc[trip_time.date==i, ['hour', 'kmh']].reset_index(drop=True).rename(columns={'kmh':var})        
                    current_day.plot(ax=ax, legend=True, lw=.6)

                ax.set_title(f'Velocidad promedio de viaje', fontsize=8)
                ax.set_xlabel('Hora', fontsize=8)
                ax.set_ylabel('Velocidad promedio\n(kmh)', fontsize=8)
                ax.set_xticks(list(range(0, 24)))

                ax.tick_params(labelsize=6);

                box = ax.get_position()
                ax.set_position([box.x0, box.y0, box.width * .87, box.height])
                ax.legend(loc='center left', bbox_to_anchor=(1, 0.5), fontsize=6)
                fig.savefig(current_path / 'Resultados_png' / f'{city}_kmh_dia_completo.png', dpi=300)
                fig.savefig(current_path / 'Resultados_pdf' / f'{city}_kmh_dia_completo.pdf', dpi=300)
                if prs == '': display(fig)

                if prs != '':
                    slide = pptx_addtitle(prs=prs, slide='',  title=title_ppt, left=0, top=0, width=24, new=True, fontsize=48)    
                    pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_tiempos_dia_completo.png',  left=.5, top=1.5, width=11)
                    pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_kmh_dia_completo.png'    ,  left=11.5, top=1.5, width=11)

                    pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_indicadores_dia_completo.png',  left=3, top=7.5, width=14)
                else:
                    display(indicators_all.set_index('Detalle día').T)

        else:
            print('No hay datos para procesar indicadores')
            indicators_all = pd.DataFrame([])
    else:
            print('No hay datos para procesar indicadores')
            indicators_all = pd.DataFrame([])
    
    return indicators_all
        
def print_density_nse(hexs, 
                      population = '', 
                      k = 4, 
                      nse = 'NSE_5', 
                      current_path = Path(), 
                      city = '',
                      prs = '',
                      title_ppt=''):
    
    utils.create_result_dir(current_path)
    
    hexs['density_ha'] = round(hexs[population] / (hexs['area_m2']/10000),1).round().astype(int)

    bins = [0] + mapclassify.NaturalBreaks(hexs['density_ha'], k=k).bins.tolist()
    bins_labels = [f'{int(bins[n])} - {int(bins[n+1])}' for n in range(0, len(bins)-1)]
    hexs['density_ha_label'] = pd.cut(hexs.density_ha, bins=bins, labels=bins_labels)

    # Find the centre of the reprojected zonas
    w,s,e,n = hexs.to_crs(3857).total_bounds
    cx = (e+w)/2 # centre point x
    cy = (n+s)/2 # centre point y

#     # For the actual plot define a crop that is tighter to minimise unused space at the edge of the map
#     crop_extent =  max(abs((e-w)/2), abs((n-s)/2))  * 1.05
#     crop_w = cx-crop_extent
#     crop_s = cy-crop_extent
#     crop_e = cx+crop_extent
#     crop_n = cy+crop_extent

    fig = Figure(figsize=(13.5,13.5), dpi=40)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)

    hexs.to_crs(3857).plot(ax=ax, column='density_ha_label', categorical = True, lw=0, alpha=.6, cmap='cividis_r', legend=True,
                          legend_kwds={'loc': 'best', 'frameon': True, 'edgecolor':'black', 'facecolor':'white', "title":'Densidad\n(Pers/ha)', 'title_fontsize':14, 'fontsize':14})

    ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution=None, attribution_size=10)

#     # Manual plot settings
#     ax.set_xlim(crop_w,crop_e)
#     ax.set_ylim(crop_s,crop_n)

    ax.set_title('Densidad Poblacional', fontsize=18)
    ax.axis('off');

    fig.tight_layout()
    fig.savefig(current_path / 'Resultados_png'  / f'{city}_density_ha.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf'  / f'{city}_density_ha.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig) 
        
    hexs['NSE_X'] = hexs[nse].replace({'1 - Alto':'Alto', '2 - Medio-Alto':'Medio-Alto', '3 - Medio':'Medio',  '4 - Medio-Bajo':'Medio-Bajo', '5 - Bajo':'Bajo'})
    if hexs['NSE_X'].dtype == 'O':
        hexs['NSE_X'] = pd.CategoricalIndex(hexs['NSE_X'], categories= ['Alto', 'Medio-Alto', 'Medio', 'Medio-Bajo', 'Bajo'])
        
    
    fig = Figure(figsize=(13.5,13.5), dpi=40)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)

    hexs.to_crs(3857).plot(ax=ax, column='NSE_X', categorical = True, lw=0, alpha=.6, cmap='YlOrRd', legend=True,
                          legend_kwds={'loc': 'best', 'frameon': True, 'edgecolor':'black', 'facecolor':'white', "title":'Nivel Socioeconómico', 'title_fontsize':14, 'fontsize':14})
    
    ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution=None, attribution_size=10)
    
    del hexs['NSE_X']
    
#     # Manual plot settings
#     ax.set_xlim(crop_w,crop_e)
#     ax.set_ylim(crop_s,crop_n)
    
    
    ax.set_title('Nivel Socioeconómico', fontsize=18)
    ax.axis('off');
    fig.tight_layout()
    fig.savefig(current_path / 'Resultados_png'  / f'{city}_NSE.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf'  / f'{city}_NSE.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig) 
        
    if prs != '':
        slide = pptx_addtitle(prs=prs, slide='',  title=title_ppt, left=0, top=0, width=24, new=True, fontsize=48)    
        pptx_addpic(prs=prs, slide=slide, img_path= current_path / 'Resultados_png'  / f'{city}_NSE.png',          left=.5, top=1.5, width=11)
        pptx_addpic(prs=prs, slide=slide, img_path= current_path / 'Resultados_png'  / f'{city}_density_ha.png',  left=12.5, top=1.5, width=11)
    
    
        
def calculate_avg_time_distance(hexs, 
                                od_matrix,                                 
                                population=''):
    
    indicators_vars = [i for i in od_matrix.columns if ('time' not in i)&(('transit_' in i)|('driving_' in i)|('osm_' in i)|('walking_' in i)|('bicycling_' in i))]

    common_vars = ['area_m2', population, 'PCA_1', 'NSE_5', 'NSE_3', 'trip_datetime']

    od_matrix_agg = od_matrix.groupby('hex')[common_vars].first().reset_index()

    for var in indicators_vars:
        val = od_matrix[od_matrix[var].notna()].hex.value_counts().value_counts(normalize=True).reset_index().rename(columns={'index':'qty'})
        if len(val) > 0:
            val['hex_agg'] = val.hex.cumsum()
            for i, row in val.iterrows():
                if row.hex_agg > .85:
                    qty_val = row.qty
                    break

            val = od_matrix[od_matrix[var].notna()].hex.value_counts().reset_index().rename(columns={'index':'hex', 'hex':'qty_val'})
            val = val[val.qty_val >= qty_val]
            od_matrix_val = od_matrix.loc[od_matrix.hex.isin(val.hex.tolist()), ['hex', var, 'weight']].reset_index(drop=True)

            od_matrix_val = od_matrix_val.groupby('hex').apply(lambda x: np.average(x[var], weights=x['weight'])).round(2).reset_index().rename(columns={0:var})
            od_matrix_agg = od_matrix_agg.merge(od_matrix_val, how='left', on='hex')


    downtown = od_matrix.loc[od_matrix.weight == od_matrix.weight.max(), ['hex']+indicators_vars]
    for i in downtown.columns:
        if i != 'hex':
            downtown = downtown.rename(columns = {i: i+'_downtown'})

    indicators_vars = indicators_vars + downtown.columns[1:].tolist()

    od_matrix_agg = od_matrix_agg.merge(downtown, how='left', on='hex')

    od_matrix_agg = hexs[['hex', 'geometry']].merge(od_matrix_agg, on='hex')
    
    return od_matrix_agg

def distribution_travel_times(od_matrix,
                              current_path = Path(), 
                              city='',
                              title_ppt = '',
                              prs=''):
    
    if ('transit_duration' in od_matrix.columns)&('driving_duration_in_traffic' in od_matrix.columns):
    
        with sns.axes_style('white', {"axes.facecolor": "#d4dadc", 'figure.facecolor': "#d4dadc"}):
            # Create figure, canvas and axis
            fig = Figure(figsize=(6.5,4), dpi=200)
            canvas = FigureCanvas(fig)
            ax = fig.add_subplot(111)
            # Add the seaborn plots
            sns.distplot(od_matrix.transit_duration, label="Transporte público", color='#5ab4ac', ax=ax)
            sns.distplot(od_matrix.driving_duration_in_traffic, label="Automóvil", color='#d8b365', ax=ax)
            # Manually adjust ax settings
            ax.set_title("Comparación de tiempos de viaje", fontsize=18, weight='bold')
            ax.set_xlabel("Tiempos de viaje (minutos)", fontsize=14)    
            ax.set_ylabel("") 
            ax.legend(fontsize=14, edgecolor='white')
        #     ax.set_xticks([0, 100, 200])
            ax.set_yticks([])
            ax.tick_params(labelsize=14)
            # Save the plot, specifying the figure background colour
            fig.tight_layout()
            fig.savefig(current_path / 'Resultados_png' / f'{city}_distplot.png', facecolor='#d4dadc', dpi=300)
            fig.savefig(current_path / 'Resultados_pdf' / f'{city}_distplot.pdf', facecolor='#d4dadc', dpi=300)
        
        if prs != '':
            slide = pptx_addtitle(prs=prs, slide='',  title=title_ppt, left=0, top=0, width=24, new=True, fontsize=48)    
            pptx_addpic(prs=prs, slide=slide, img_path= current_path / 'Resultados_png' / f'{city}_distplot.png',  left=4, top=3, width=15)



def print_graphs(   od_matrix_new,
                    od_matrix_new_w,
                    var,                     
                    colors_dict = '',                    
                    k = 5,
                    bins = [],
                    extend='neither',
                    population = '',
                    equipment_type = '',
                    equipment_type_title = '',
                    alpha=.6,
                    current_path = Path(),
                    city = '',
                    prs = '',
                    title_ppt='',
                    showfliers=False):
    
    '''
    extend='neither','min'
    '''
    warnings.filterwarnings("ignore")

    colors = {  'distance_osm_drive': 'YlGn',
                'distance': 'YlGn',
                'transit_distance': 'YlGn',
                'transit_walking_distance': 'autumn_r',
                'transit_transit_distance': 'YlGn',
                'transit_walking_distance_origin': 'autumn_r',
                'transit_distance_downtown': 'YlGnBu',
                'transit_walking_distance_downtown': 'YlGnBu',
                'transit_transit_distance_downtown': 'YlGnBu',
                'transit_walking_distance_origin_downtown': 'YlGnBu',
                'driving_distance': 'YlGn',
                'bicycling_distance': 'YlGn',
                'distance_osm_walk': 'YlGn',
                'distance_osm_walk_downtown': 'YlGnBu',
                'walking_distance': 'YlGn',
                'distance_osm_drive_downtown': 'YlGnBu',
                'driving_distance_downtown': 'YlGnBu',
                'transit_duration': 'PuBuGn',
                'duration': 'PuBuGn',
                'transit_walking_duration': 'PuBuGn',
                'transit_walking_duration_origin': 'PuBuGn',
                'transit_duration_downtown': 'YlOrBr',
                'transit_walking_duration_downtown': 'YlOrBr',
                'transit_walking_duration_origin_downtown': 'YlOrBr',
                'driving_duration_in_traffic': 'bone_r',
                'driving_duration_in_traffic_downtown': 'pink_r',
                'walking_duration': 'PuBuGn',
                'bicycling_duration': 'PuBuGn',
                'transit_transit_steps': 'YlOrRd',
                'transit_transit_steps_downtown': 'YlOrRd' ,
                'transit_transit_steps_downtown': 'Greens' }
    
    try:
        cmap = colors_dict[var]
    except:
        try:
            cmap = colors[var]
        except:
            cmap = 'YlGn'
    
    desc = ''
    label = ''
    label_short = ''
    legtitle = ''
    if 'distance' in var: desc += 'Distancias'
    if 'duration' in var: desc += 'Tiempos'
    if 'steps' in var: desc += 'Etapas'
    if 'transit' in var: desc += ' en transporte público'
    if 'driving' in var: desc += ' en automóvil'
    if 'walking' in var: desc += ' caminando'
    if 'bicycling' in var: desc += ' en bicicleta'
    if 'transit_walking_distance' in var: desc = 'Distancias caminadas (viajes en transporte público)'
    if 'transit_walking_distance_origin' in var: desc = 'Cobertura de transporte público'
    if 'transit_walking_duration' in var: desc = 'Tiempo caminando a la parada'
        
    if 'green_area_m2' in var: desc = 'Áreas verdes per capita (m2) en '+var[var.find('_in')+4:] + 'mts'
    
    if 'downtown' in var: desc += '\n(viajes al centro)'

    if 'distance' in var: label = 'kms'
    if 'duration' in var: label = 'minutos'
    if 'steps' in var: label = 'steps'
    if 'walking_distance' in var: label = 'mts'
    if 'green_area_m2' in var: label = 'm2'

    if 'distance' in var: label_short = 'Distancia (kms)'
    if 'duration' in var: label_short = 'Tiempos de viaje (minutos)'
    if 'steps' in var: label_short = 'Modos utilizados'
    if 'walking_distance' in var: label_short = 'Distancia (mts)'
    if 'green_area_m2' in var: label_short = 'Áreas Verdes per cápita (m2)'

    if 'Distancias' in desc:
        legtitle = 'Distancias (kms)'
    if 'Tiempos' in desc:
        legtitle = 'Tiempos (min)'
    if 'Etapas' in desc:
        legtitle = 'Etapas'
    if 'áreas verdes' in desc:
        legtitle = 'Áreas Verdes (ha)'
     
    od_matrix_new = od_matrix_new[od_matrix_new[var].notna()].copy()
    od_matrix_new_w = od_matrix_new_w[od_matrix_new_w[var].notna()].copy()
        
    if equipment_type_title != 'x':        
        if var == 'distance':
            legtitle = 'Distancias (mts)'
            od_matrix_new['distance'] = od_matrix_new['distance'] * 1000
            od_matrix_new_w['distance'] = od_matrix_new_w['distance'] * 1000
            desc = desc+' '+equipment_type_title    
        if var == 'duration':
            desc = 'Tiempos de viaje - ' + equipment_type_title + '\nTransporte público o caminando'
        if 'qty_est' in var:
            desc = 'Establecimientos en '+var[var.find('est_')+4:]+' - ' + equipment_type_title      
    
    if 'walking_distance' in var:
        od_matrix_new[var] = od_matrix_new[var] * 1000
        od_matrix_new_w[var] = od_matrix_new_w[var] * 1000
        
        
    sfile = ''
    for i in od_matrix_new[equipment_type]:
        sfile += od_matrix_new[i].unique()[0]+'_'
    sfile = sfile[:-1]
    
    utils.create_result_dir(current_path)  
    
    if len(bins)==0:
        if 'steps' in var:
            bins = mapclassify.Quantiles(od_matrix_new[var], k=3).bins.tolist()
            

        elif ('transit_duration' in var)|('driving_duration' in var):            

            min = int(od_matrix_new[var].min())
            bins=[]
            for n in range(0, 4):                
                min += 15        
                bins += [min]
            bins = bins + [int(od_matrix_new[var].max())]            
        elif ('transit_walking_distance' in var):
            bins = [500, 1000, 1500, 2000]
            
        elif ('green_area_m2' in var):
            bins = [1, 3, 6, 9, 12]
            
        else:
            bins = mapclassify.Quantiles(od_matrix_new[var], k=k).bins.tolist()            

    

    bins_labels = []
    lower = int(od_matrix_new[var].min())
    for i in bins:        
        
        if 'steps' not in var: 
            upper = int(round(i))            
        else:            
            upper = round(i, 2)
            
        bins_labels = bins_labels + [f'{lower} - {upper}']
        lower = upper
    if bins[len(bins)-1] < int(od_matrix_new[var].max()):
        bins = bins+[int(od_matrix_new[var].max())]
        bins_labels = bins_labels+[f'{upper} - {int(od_matrix_new[var].max())}']

    
    ### Imprimo Mapa
    
    bins_mapa = bins.copy()
    if extend=='neither':
        if bins_mapa[0] > int(od_matrix_new[var].min()):
            bins_mapa = [int(od_matrix_new[var].min())] + bins_mapa

        if bins_mapa[len(bins_mapa)-1] < int(od_matrix_new[var].max()):
            bins_mapa = bins_mapa + [int(od_matrix_new[var].max())]

    # Find the centre of the reprojected zonas
    w,s,e,n = od_matrix_new.to_crs(3857).total_bounds
    cx = (e+w)/2 # centre point x
    cy = (n+s)/2 # centre point y

#     # For the actual plot define a crop that is tighter to minimise unused space at the edge of the map
#     crop_extent =  max(abs((e-w)/2), abs((n-s)/2)) * 1.05
#     crop_w = cx-crop_extent
#     crop_s = cy-crop_extent
#     crop_e = cx+crop_extent
#     crop_n = cy+crop_extent
    
    fig = Figure(figsize=(13.5,13.5), dpi=35)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)
    
    plt.rcParams.update({"axes.facecolor": '#d4dadc', 'figure.facecolor': '#d4dadc'})   
    
    od_matrix_new.to_crs(3857).plot(ax=ax, 
                                    column=var, 
                                    alpha=alpha, 
                                    lw=0, 
                                    categorical = True, 
                                    legend=False, 
                                    cmap = cmap, 
                                    scheme="User_Defined", classification_kwds=dict(bins=bins_mapa)
                                    )

    ax.set_title(desc, fontsize=16, weight='bold')
    ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution=None, attribution_size=10)

#     # Manual plot settings
#     ax.set_xlim(crop_w,crop_e)
#     ax.set_ylim(crop_s,crop_n)

    ax.axis('off')


    # colorbar
    cax_position = "top"
    cax_size = "1.7%" #tamaño
    cax_pad = "-4%" #Altura
    cax_aspect = 0.050 #ancho

    cmap = mpl.cm.get_cmap(cmap)
    norm = mpl.colors.BoundaryNorm(bins_mapa, cmap.N, extend=extend)

    # create an axes on the side of ax. 
    divider = make_axes_locatable(ax)
    cax = divider.append_axes(cax_position, size=cax_size, pad=cax_pad, aspect=cax_aspect)
    mpl.colorbar.ColorbarBase(cax, 
                              cmap=cmap, 
                              norm=norm, 
                              spacing='uniform', 
                              orientation='horizontal', 
                              extend=extend,
                              extendfrac=.2,
                              label=label,
                              format=tick.FormatStrFormatter('%.0f'), 
                              alpha=alpha)

    fig.tight_layout()

    fig.savefig(current_path / 'Resultados_png' / f'{city}_{var}_{sfile}_map.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf' / f'{city}_{var}_{sfile}_map.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig)

    for x in range(0, len(od_matrix_new_w['NSE_5'].unique())):
        od_matrix_new_w['NSE_5'] = od_matrix_new_w['NSE_5'].str.replace(f'{x+1} - ', '')

    ### Imprimo boxplot
    
    # Create figure, canvas, axis
    fig = Figure(figsize=(7,4.5), dpi=70)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)

    plt.rcParams.update({"axes.facecolor": '#d4dadc', 'figure.facecolor': '#d4dadc'})   

    sns.boxplot(x='NSE_5', y=var, data=od_matrix_new_w, palette='PuOr', order=['Bajo', 'Medio-Bajo', 'Medio', 'Medio-Alto', 'Alto'], ax=ax, showfliers=showfliers) 

    ax.set_title(label=desc, fontsize=13, weight='bold')
    ax.tick_params(labelsize=12)
    ax.set_xlabel('Nivel socioeconómico', fontsize=12)
    ax.set_ylabel(ylabel=label_short, fontsize=12)

    fig.tight_layout()

    fig.savefig(current_path / 'Resultados_png' / f'{city}_{var}_{sfile}_boxplot.png', facecolor='#d4dadc', dpi=300);
    fig.savefig(current_path / 'Resultados_pdf' / f'{city}_{var}_{sfile}_boxplot.pdf', facecolor='#d4dadc', dpi=300);
    if prs == '': display(fig)
            
    
    ### Imprimo gráfico de porcentajes de la población
    


    od_matrix_new_w['bins_w'] = pd.cut(od_matrix_new_w[var], [0]+bins, labels = bins_labels)
    
    od_matrix_new_w.loc[(od_matrix_new_w['bins_w'].isna())&(od_matrix_new_w[var]<bins[0]), 'bins_w'] = bins_labels[0]
    
    od_matrix_new_w['NSE_w'] = od_matrix_new_w['NSE_5'].replace({'Alto': '5 - Alto', 'Medio-Alto':'4 - Medio-Alto', 'Medio':'3 - Medio', 'Medio-Bajo':'2 - Medio-Bajo', 'Bajo':'1 - Bajo'})
    
    # Create figure, canvas and axis
    fig = Figure(figsize=(7,4.5), dpi=70)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)
    
    plt.rcParams.update({"axes.facecolor": '#d4dadc', 'figure.facecolor': '#d4dadc'})   

    dfw = (pd.crosstab(index=od_matrix_new_w['NSE_w'],
                        columns=od_matrix_new_w['bins_w'], 
                        values=od_matrix_new_w[population], 
                        aggfunc='sum', 
                        normalize='index').fillna(0)) * 100
    dfw.plot(kind='bar', 
              stacked=True, 
              cmap=cmap,
              ax=ax)

    # Manually adjust ax settings
    ax.set_title(desc, fontsize=13, weight='bold')

    ax.tick_params(labelsize=11)
    ax.xaxis.set_tick_params(rotation=0)
    ax.set_xlabel('Nivel Socioeconómico', fontsize=11, rotation='horizontal')
    ax.set_ylabel('Porcentaje de la población', fontsize=11)

    box = ax.get_position()
    ax.set_position([box.x0, box.y0, box.width * 0.75, box.height])
    # Put a legend to the right of the current axis
    ax.tick_params(axis='both', which='major', labelsize=8)
    ax.legend(loc='center left', bbox_to_anchor=(1, 0.5), fontsize=8, title=legtitle, title_fontsize=8, frameon='white')

    # Save the figure, specifying the figure background colour
    fig.subplots_adjust(right=0.75, bottom=0.4)
    fig.tight_layout()

    labels = [i.get_text() for i in ax.get_xticklabels()]
    labels = [i[4:] for i in labels]
    ax.set_xticklabels(labels)
    
    fig.tight_layout()

    fig.savefig(current_path / 'Resultados_png' / f'{city}_{var}_{sfile}_bar.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf' / f'{city}_{var}_{sfile}_bar.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig)

    if prs != '':
        slide = pptx_addtitle(prs=prs, slide='',  title=title_ppt, left=0, top=0, width=24, new=True, fontsize=48)    
        pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_{var}_map.png',  left=0, top=1.3, width=12)
        pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_{var}_boxplot.png', left=14, top=1, width=9)
        pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_{var}_bar.png', left=14.5, top=7, width=9)


def print_time_distance(od_matrix, 
                        hexs='',
                        indicators_vars='', 
                        equipment_type='',
                        colors_dict = '',
                        current_path=Path(), 
                        population = '',
                        k = 5,
                        bins=[],
                        extend='neither',
                        alpha=.6,
                        city='', 
                        prs='',
                        title_ppt='',                        
                        showfliers=False):
    
    if len(indicators_vars) == 0:         
        indicators_vars=od_matrix.columns
    
        indicators_vars = [i for i in od_matrix.columns if (i in indicators_vars)&
                                                          (('transit_' in i)|
                                                           ('driving_' in i)|
                                                           ('osm_' in i)|
                                                           ('walking_' in i)|
                                                           ('bicycling_' in i)|
                                                           ('distance' == i)|
                                                           ('distance' == i)|
                                                           ('duration' == i)|
                                                           ('green_area_m2' in i))]

        
    if len(equipment_type) == 0:
        od_matrix['equipment_type_tmp'] = 'x'
        equipment_type = ['equipment_type_tmp']

    if type(equipment_type) == str:
        equipment_type = [equipment_type]

    for _, i in od_matrix.groupby(equipment_type).size().reset_index().iterrows():   
        
        if (equipment_type != ['equipment_type_tmp'])&(len(colors_dict)==0):            
            txt = ''
            for x in i[equipment_type].values.tolist(): txt += x +' '

            if 'Inicial' in txt:
                colors_dict = {'distance':'Reds', 'duration':'summer_r'}
            if 'Primaria' in txt:
                colors_dict = {'distance':'Reds', 'duration':'pink_r'}
            if 'Secundaria' in txt:
                colors_dict = {'distance':'Reds', 'duration':'hot_r'}
            if 'Atención Primaria' in txt:
                colors_dict = {'distance':'Reds', 'duration':'gist_heat_r'}
            if 'Hospitales' in txt:
                colors_dict = {'distance':'Reds', 'duration':'copper_r'}
        
        df = od_matrix[(od_matrix[equipment_type] == i[equipment_type]).all(axis=1)].copy()
        
        if len(hexs) > 0:
            df = hexs[['hex', 'geometry']].merge(df)
    
        df_w = utils.reindex_df(df, weight_col = population, div=10)
        
        equipment_type_title = i[equipment_type].values[0]
        for x in i[equipment_type].values[1:].tolist():
            equipment_type_title += ' - '+x
        
        for var in indicators_vars:
            
            print_graphs(df,
                         df_w, 
                         var,
                         colors_dict = colors_dict,
                         population = population,
                         equipment_type = equipment_type,
                         equipment_type_title = equipment_type_title,
                         k = k,
                         bins = bins,
                         extend=extend,
                         alpha=alpha,
                         current_path = current_path,
                         city = city,
                         prs = prs,
                         title_ppt=title_ppt,
                         showfliers=showfliers)
                         
                         
def create_pptx(hexs, 
                od_matrix_all_day='',
                od_matrix='',                
                od_establecimientos='',
                hexs_green_space = '',
                population='',
                equipment_type='',
                current_path = Path(), 
                city='',
                showfliers=False,
                title_ppt = ''):

    warnings.filterwarnings("ignore")
    
    prs = Presentation()
    prs.slide_height = Inches(13.5)
    prs.slide_width = Inches(24)

    file_pptx = current_path / 'Resultados_pptx' / f'{city}_Accesibilidad.pptx'

    print('Densidad y nivel socioeconómico')
    print_density_nse(hexs, 
                      population=population,
                      current_path = current_path, 
                      city=city,
                      prs=prs,
                      title_ppt=title_ppt)

    if len(od_matrix_all_day) > 0:
        print('Índicadores de día completo')
        indicators_all = indicators_all_day(od_matrix_all_day, current_path, city=city, prs=prs, title_ppt=title_ppt)

    if len(od_matrix) > 0:
        
        distribution_travel_times(od_matrix,
                         current_path = current_path, 
                         city=city,
                         title_ppt = title_ppt,
                         prs=prs)
        
        
        indicators_vars = ['transit_duration',
                           'driving_duration_in_traffic',
                           'transit_walking_distance', 
                           'transit_walking_distance_origin',                    
                           'walking_duration', 
                           'bicycling_duration']

        od_matrix_avg = calculate_avg_time_distance(hexs, 
                                                    od_matrix,
                                                    population=population)

        indicators_vars = indicators_vars + ['transit_duration_downtown', 
                                             'driving_duration_in_traffic_downtown',
                                             'walking_duration_downtown', 
                                             'bicycling_duration_downtown']

        try:
            print('Isocronas de tiempos y distancias')
            print_time_distance(od_matrix_avg, 
                                hexs=hexs,
                                indicators_vars=indicators_vars, 
                                current_path=current_path, 
                                population=population,
                                city=city, 
                                prs=prs,
                                title_ppt=title_ppt,
                                showfliers=showfliers)
        except:
            print('Error imprimiendo mapas de tiempos y distancias')
            pass


    if len(od_establecimientos) > 0:
        print('Establecimientos')

        print_time_distance(od_establecimientos, 
                            hexs=hexs,
                            population=population,
                            indicators_vars=['duration'],
                            colors_dict={'distance':'Reds', 'duration':'Blues'},
                            equipment_type = equipment_type,
                            current_path=current_path, 
                            city=city,
                            prs=prs,
                            title_ppt=title_ppt,
                            showfliers=showfliers)
    
    if len(hexs_green_space) > 0:
        
        green_area_var = [_ for _ in hexs_green_space.columns if 'green_area_m2' in _]
        green_area_var = green_area_var[0]
        print_time_distance(hexs_green_space, 
                            population=population,
                            indicators_vars=[green_area_var],
                            colors_dict={green_area_var:'Greens'},                    
                            current_path=current_path, 
                            city=city,
                            prs=prs,
                            title_ppt=title_ppt,
                            showfliers=showfliers)


    try:
        print('')
        prs.save(file_pptx)
        print(file_pptx)
        print('')
    except:
        print('')
        print('No se pudo guardar el archivo', file_pptx)
