import pandas as pd
from pandas._libs.lib import is_integer
import numpy as np
import geopandas as gpd
import h3
import matplotlib.pyplot as plt
import seaborn as sns

import mapclassify
import contextily as ctx

from tqdm.notebook import tqdm
tqdm.pandas()


import unidecode
import itertools
from pathlib import Path

from sklearn.cluster import DBSCAN
from sklearn import metrics
from sklearn.metrics import silhouette_samples, silhouette_score

from sklearn.preprocessing import StandardScaler, RobustScaler 
from sklearn.decomposition import PCA

import osmnx as ox
import networkx as nx
import pandana

from pandana.loaders import osm as osm_pandana

import datetime as datetime
from datetime import date
from datetime import timedelta

pd.set_option('display.max_columns', None)

import googlemaps

import pytz
from tzwhere import tzwhere 
from pytz import timezone

import warnings
warnings.filterwarnings("ignore")
warnings.filterwarnings("ignore", category=DeprecationWarning) 


from shapely.geometry import Point, Polygon

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor # ColorFormat, 
from PIL import Image, ImageDraw

import os
import glob

import matplotlib as mpl
from mpl_toolkits.axes_grid1 import make_axes_locatable
import matplotlib.pyplot as plt
import matplotlib.ticker as tick
import IPython
from IPython.display import set_matplotlib_formats


from PIL import Image # pip install Pillow
from PIL import ImageOps

import seaborn as sns
from pandas import DataFrame

from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas
from matplotlib.figure import Figure
from matplotlib import cm
from matplotlib.colors import ListedColormap, LinearSegmentedColormap

from mpl_toolkits.axes_grid1 import make_axes_locatable


from pyomu.support_h3 import support_h3 as suph3
from pyomu.support import support as sup

    
def calculate_nse(df, vars, population, var_result='PCA_1', scaler_type='Standard', q=[5, 3], order_nse='', show_map=False):
    '''
    Realiza el análisis de Componentes Principales (PCA) para un conjunto de variables en un dataframe. Este resultado se utiliza para la
    creación de una variable de nivel socioeconómico
    
    Parámetros: 
    df = dataframe que contiene las variables que se utilizan para calcular el PCA y nivel socieconómico. El dataframe puede ser de hexágonos, 
    radios censales u otra geometría con datos de los hogares o población.
    vars = variables que se van a utilizar para el calculo de PCA
    population = variable población.
    var_result= variable donde se alojará el resultado del PCA. Por defecto será 'PCA_1'
    scaler_type=Tipo de normalización para el PCA. Por defecto será 'Standard'
    q = Cantidad de agrupaciones para el nivel socioeconómico resultado. Por defecto serán quintiles y terciles = [5, 3]
    show_map = Muestra mapa resultado, por defecto es False
    
    Salida: DataFrame similar al de entrada con nuevas variables que contienen el PCA y los grupos solicitados (i.e. 'PCA_1', 'NSE_5', 'NSE_3')
    '''

    if type(var_result) == str:
        var_result = [var_result]
    
    for i in var_result:
        if i in df.columns:
            df.drop([i], axis=1, inplace=True)
    
    data_1 = df[vars]
    if scaler_type!='None':
        if scaler_type=='Standard':        
            scaler = StandardScaler()
        if scaler_type=='Robust':        
            scaler = RobustScaler()
        
        data_1 = scaler.fit_transform(data_1)

    # Importamos la clase PCA del modulo decomposition de la librería Sklearn.

    # Instanciamos la clase pidiendo que conserve la cantidad de componentes requerida
    pca_2cp = PCA(n_components=len(var_result), svd_solver='full')

    # Con el método fit() calculammos los componentes principales
    principalComponents = pca_2cp.fit_transform(data_1)
    
    principalDf = pd.DataFrame(data = principalComponents, 
                               columns = var_result)

    df_result=pd.concat([df, principalDf], axis=1)

    print('variance ratio', round(pca_2cp.explained_variance_ratio_[0], 2))
    print('(% de la variancia explicada por el componente 1)')
    print('')
    
    
    # Calcula el NSE ponderado por la población (puedo calcular NSE para varios q según parámetros)
    if type(q) == int: 
        q = [q]
        
    for i in q:    
        
        df_result[f'NSE_{i}'] = sup.weighted_qcut(df_result[var_result[0]], df_result[population], i, labels=False)
        df_result[f'NSE_{i}'] = df_result[f'NSE_{i}'] + 1
    

    if len(order_nse) > 0:
        for i in range(0, len(order_nse)):
            
            lst = [x+1 for x in range(0, q[i])]
            ord = order_nse[i]        
            ord = [f'{x+1} - {ord[x]}' for x in range(0, len(ord))]        
            df_result[f'NSE_{q[i]}'] = df_result[f'NSE_{q[i]}'].replace(lst, ord)
            
    
    if show_map:
        
        df_result['NSE_X'] = df_result[f'NSE_{i}'].replace({'1 - Alto':'Alto', '2 - Medio-Alto':'Medio-Alto', '3 - Medio':'Medio',  '4 - Medio-Bajo':'Medio-Bajo', '5 - Bajo':'Bajo'})
        df_result['NSE_X'] = pd.CategoricalIndex(df_result['NSE_X'], categories= ['Alto', 'Medio-Alto', 'Medio', 'Medio-Bajo', 'Bajo'])
        
        fig, ax = plt.subplots(dpi=150, figsize=(5, 5))
        df_result.to_crs(3857).plot(column='NSE_X', 
                                    ax=ax, 
                                    cmap='Reds',
                                    categorical=True,
                                    legend=True,                            
                                    legend_kwds={'loc': 'best', 'frameon': True, 'edgecolor':'white', 'facecolor':None, "title":'NSE', 'title_fontsize':8, 'fontsize':7})

        ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=3)
        ax.axis('off');
        del df_result['NSE_X']
    
    return df_result
    
    


def distribute_population(gdf, 
                          id_gdf, 
                          hexs, 
                          id_hexs, 
                          population, 
                          pca, 
                          crs, 
                          q=[5, 3], 
                          order_nse='',
                          verify_overlay=True,
                          verify_overlay_df = '',
                          verify_overlay_df_id = '',
                          tags_osm = {'leisure':True, 'aeroway':True, 'natural':True},
                          show_map=False):

    
    
    '''
    Distribuye la población y el resultado del PCA de una capa geográfica principal a una segunda capa geográfica. Por ejemplo,
    permite distribuir la población y el PCA calculado para una capa geográfica censal en una capa geográfica de hexágonos.
    La población la distribuye según el área ocupada de cada polígono en relación a los polígonos de la segunda capa y el PCA para la 
    segunda capa geográfica lo calcula como un promedio ponderado (usando la variable población) de los PCAs de la capa principal.
    A su vez, permite verificar si existe superposición con otra capa geográfica si es posible identificar areas donde no hay población.
    Por ejemplo, esto sirve si se pueden identificar espacios verdes o públicos para asignar mejor la población en áreas donde no existen este tipo de
    equipamientos. De no existir esta capa geográfica, se pueden obtener y utilizar una capa obtenida de Open Street Maps. Por defecto obtiene los elementos
    relacionados con recreación, aeronavegación y espacios naturales.
    
    gdf = Capa geográfica principal (por ejemplo, una capa de radios censales)
    id_gdf = ID de la capa principal
    hexs = Capa geográfica secundaria (por ejemplo, una capa de hexágonos)
    id_hexs = ID de la capa de secundaria
    population = variable donde está la población de cada polígono en la capa principal
    pca = variable donde está el pca de cada polígono en la capa principal
    crs = Proyección correspondiente a la capa en metros (es importante que no sea un proyección de en grados)
    q=[5, 3] = Resultado de las variables de nivel socioeconómico que se quiere contruir. Por defecto se construyen quintiles y terciles. q = [5, 3]
    verify_overlay = Verifica que exista superposición con otra capa geográfica donde se puede identificar que no hay población (por ejemplo parques o espacios naturales)
    verify_overlay_df = GeoDataFrame con la capa geográfica que se quiere verificar la superposición
    verify_overlay_df_id = ID del GeoDataFrame con la capa geográfica que se quiere verificar la superposición
    tags_osm = Si se quiere obtener la información de superposición de Open Street Mapas. Por defecto trae los siguientes elementos: 
    tags_osm = {'leisure':True, 'aeroway':True, 'natural':True},
    show_map=Permite mostrar un mapa con el resultado del nuevo NSE calculado para la capa geográfica secundaria.
    
    Salida: Devuelve la capa geográfica secundaria con la distribución de las variables de población, PCA y NSE calculadas teniendo en cuenta la capa geográfica principal 
    '''
   
    shape = gpd.overlay(hexs[[id_hexs, 'geometry']], gdf[[id_gdf, population, pca, 'geometry']], how='intersection')

    shape['area_interception'] = shape.to_crs(crs).area

    if verify_overlay:
        if len(verify_overlay_df)==0:
            verify_overlay_df_id = 'osmid'
            verify_overlay_df = bring_osm(gdf, tags = {'leisure':True, 'aeroway':True, 'natural':True}, list_types = ['Polygon', 'MultiPolygon'] )

    shape_space_not_available = gpd.overlay(shape[[id_hexs, id_gdf, 'geometry']], verify_overlay_df[['osmid', 'geometry']], how='intersection')
    shape_space_not_available = shape_space_not_available[[id_hexs, id_gdf, 'geometry']].dissolve(by=[id_hexs, id_gdf]).reset_index()
    shape_space_not_available['area_not_available'] = shape_space_not_available.to_crs(crs).area


    shape = shape.merge(shape_space_not_available[[id_hexs, id_gdf, 'area_not_available']], on=[id_hexs, id_gdf], how='left').fillna(0)

    shape['area_interception_result'] = (shape['area_interception'] - shape['area_not_available']) 
    
    shape['area_gdf_result'] = shape.groupby(id_gdf).area_interception_result.transform(sum)

    shape['distribute_population'] = shape['area_interception_result'] / shape['area_gdf_result']

    shape[population] = (shape[population] * shape['distribute_population']).round()
    

    df_result = shape.groupby(id_hexs, as_index=False)[population].sum().merge(
                        shape[shape[population]>0].groupby(id_hexs).apply(lambda x: np.average(x[pca], 
                                                                                               weights=x[population])).reset_index().rename(columns={0:pca}), 
                                                                                               how='left' )

    
    df_result = hexs[[id_hexs, 'geometry']].merge(df_result, how='left')
    df_result['area_m2'] = df_result.to_crs(crs).area.round()
    df_result['density_ha'] = round(df_result[population] / (df_result['area_m2']/10000),1)
    df_result = df_result[[id_hexs, 'area_m2', population, pca, 'geometry']]
    
    # Calcula el NSE ponderado por la población (puedo calcular NSE para varios q según parámetros)

    df_result = df_result[df_result[pca].notna()].reset_index(drop=True)

    if type(q) == int: 
        q = [q]
        
    for i in q:            
        df_result[f'NSE_{i}'] = sup.weighted_qcut(df_result[pca], df_result[population], i, labels=False)
        df_result[f'NSE_{i}'] = df_result[f'NSE_{i}'] + 1
    
    
    for i in range(0, len(order_nse)):            
        lst = [x+1 for x in range(0, q[i])]      
        ord = order_nse[i]        
        ord = [f'{x+1} - {ord[x]}' for x in range(0, len(ord))]        
        df_result[f'NSE_{q[i]}'] = df_result[f'NSE_{q[i]}'].replace(lst, ord)
    
    
    if show_map:
        
        
        i = q[0]
        
        fig = plt.figure(figsize=(15,15), dpi=100)
        
        sns.set_style("white")
        
        fig.suptitle('NSE', fontsize=16)
        
        ax = fig.add_subplot(2,2,1)
        
        gdf.to_crs(3857).plot(column=f'NSE_{i}', 
                                    ax=ax, 
                                    cmap='Reds',
                                    categorical=True,
                                    legend=True,                            
                                    legend_kwds={'loc': 'best', 
                                                 'frameon': True, 
                                                 'edgecolor':'white', 
                                                 'facecolor':None, 
                                                 'title':'NSE', 
                                                 'title_fontsize':8, 
                                                 'fontsize':7})
        
        ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=5)
        
        ax.axis('off');

        ax = fig.add_subplot(2,2,2)
        
        
        df_result['NSE_X'] = df_result[f'NSE_{i}'].replace({'1 - Alto':'Alto', '2 - Medio-Alto':'Medio-Alto', '3 - Medio':'Medio',  '4 - Medio-Bajo':'Medio-Bajo', '5 - Bajo':'Bajo'})
        df_result['NSE_X'] = pd.CategoricalIndex(df_result['NSE_X'], categories= ['Alto', 'Medio-Alto', 'Medio', 'Medio-Bajo', 'Bajo'])
                
        df_result.to_crs(3857).plot(column=f'NSE_X', 
                                    ax=ax, 
                                    cmap='Reds',
                                    categorical=True,
                                    legend=True,                            
                                    legend_kwds={'loc': 'best', 
                                                 'frameon': True, 
                                                 'edgecolor':'white', 
                                                 'facecolor':None, 
                                                 'title':'NSE', 
                                                 'title_fontsize':8, 
                                                 'fontsize':7})
        
        ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=5)
        
        ax.axis('off');
        del df_result['NSE_X']
        
    
    return df_result




def assign_weights(amenities):
    '''
    Asigna ponderación a los objetos espaciales de Open Street Maps. Esta ponderación se va a utilizar para la clusterización e identificación de principales atracctores de actividad.
    Parámetros
        amenities = GeoDataFrame obtenido de OSMNX con objetos espaciales de Open Street Maps.

    Salida: agrega variable weight con la ponderación asignada a cada objeto espacial.
    '''
    
    
    cat0 = ['college',
            'university',
            'hospital',
            'clinic',        
            'bank',
            'bus_station',
            'ferry_terminal',
            'mall',
            'centro comercial',
            'multiple uso comercial',
            'juzgados',
            'courthouse',
            'public_building',
            'government',
            'register_office',
            'defensa'
            ]

    cat1 = ['education',
            'school',
            'primaria',
            'library',
            'laboratory',        
            'factory',
            'community_hall',
            'community_centre',
            'post_office'
            ]

    cat2 = [
            'police',
            'Prefectura',
            'first_aid',        
            'supermarket',
            ]


    cat3 = ['auditorium',
            'arts',
            'theatre',
            'cinema',
            'exhibition_centre',
            'comunity_center',
            'conference_center',
            'convention_centre',
            'music_venue',
            'comercio',
            'store',
            'shop',
            'comercial',
            'boutique',      
            'childcare',
            'kindergarten',
            'jardin infantil',
            'church',
            'place_of_worship',
            'coworking_space',
            'doctors',      
            'embassy',
            'emergency_service',
            'entertainment',   
            'fast_food',
            'cafeteria',
            'restaurant',
            'bar',
            'gambling',
            'gimnasium',
            'gym',
            'social_club',
            'sport_centre',
            'laboratorio',
            'odontologia',
            'office',
            'oficina',
            'payment_centre'
            'peluqueria'
            'pharmacy',
            'recreation',        
            'pharmacy'        
    ]

    tags_weighted = pd.DataFrame()
    w = [100, 50, 20, 10]
    n = 0
    for i in [cat0, cat1, cat2, cat3]:
        tags_ = pd.DataFrame(i, columns=['tag_type'])
        tags_['weight'] = w[n]
        tags_weighted = pd.concat([tags_weighted, 
                                   tags_], ignore_index=True)
        n += 1

    amenities['tag_type'] = amenities.tag_type.apply(sup.normalize_words)

    for i in tags_weighted.sort_values('weight', ascending=True).itertuples():   
        
        amenities.loc[amenities.tag_type.str.contains(i.tag_type), 'weight'] = i.weight

    amenities['weight'] = amenities['weight'].fillna(1)
    
    amenities = amenities[amenities.weight > 0].reset_index(drop=True)

    return amenities.sort_values('tag_type').reset_index(drop=True)


def activity_density(amenities, city_crs, cantidad_clusters = 15, show_map = False): 
    '''
    Teniendo en consideración un GeoDataFrame de equipamientos se identifican los clusters de alta densidad de establecimientos.
    Se sugiere la obtención del GeoDataFrame de Open Stree Maps y se la asigna una ponderación a cada establecimiento (ver funciones bring_osm y assign_weights). 
    Este GeoDataFrame puede ser reemplazo por alguna otra fuente. De no existir la variable de ponderación (weight), se le asigna 1 a cada registro.
    
    Parámetros:
    amenities = GeoDataFrame con los equipamientos
    city_crs = proyección en metros de la ciudad de análisis para el cálculo de distancias entre establecimientos
    cantidad_clusters = Cantidad de clusters que se quiere obtener. Por defecto 15.
    show_map = Muestra mapa resultado del proceso con los clusters de actividad. Por defecto False
    
    Salida: Devuelve un nuevo GeoDataFrame con los principales clusters de actividad
    '''

    eps_ = [250, 500, 750, 1000]
    samples_ = list(range(500, 3000, 200)) 

    amenities['geometry'] = amenities['geometry'].representative_point()
    
    if 'weight' not in amenities.columns:
        amenities['weight'] = 1

    amenities['x'] = amenities.to_crs(city_crs).geometry.centroid.x
    amenities['y'] = amenities.to_crs(city_crs).geometry.centroid.y
    X = amenities.reindex(columns = ['x','y']).values
    W = amenities.weight

    scores = pd.DataFrame([])

    for eps, samples in tqdm(list(itertools.product(eps_, samples_)), desc='Densidad de actividad'):

        cluster_name = f'cluster_{eps}_{samples}'

        clustering = DBSCAN(eps=eps, 
                            min_samples=samples 
                            ).fit(X, sample_weight=W)

        amenities[cluster_name] = clustering.labels_

        try:
            sc = metrics.silhouette_score(X, clustering.labels_).round(2)
            exception = False
        except:
            exception = True
            sc = -1

        if len(amenities[amenities[cluster_name]>-1])>0:
            cant_clusters = len(amenities.loc[(amenities[cluster_name]>-1),cluster_name].unique())

            result = amenities[amenities[cluster_name]>-1].groupby(cluster_name, as_index=False).agg({'weight':'sum', 'x': 'mean', 'y':'mean'}).sort_values('weight', ascending=False)

            result['weight'] = result['weight'] / amenities['weight'].sum() * 100

            scores = pd.concat([scores, 
                                pd.DataFrame([[cluster_name, 
                                               eps, 
                                               samples, 
                                               sc, 
                                               cant_clusters, 
                                               result.head(1).weight.values[0], 
                                               result.tail(1).weight.values[0], 
                                               exception]], 
                                            columns=['cluster', 
                                                     'eps', 
                                                     'samples', 
                                                     'score', 
                                                     'cant_clusters', 
                                                     'max_weight', 
                                                     'min_weight', 
                                                     'exception'])], ignore_index=True)

    scores = scores[(scores.cant_clusters>=cantidad_clusters)].sort_values(['eps', 'cant_clusters']).reset_index(drop=True)

    cluster_name = scores.head(1).cluster.values[0]

    result = amenities[amenities[cluster_name]>-1].groupby(cluster_name, as_index=False).agg({'weight':'sum'}).sort_values('weight', ascending=False)
    result = result.merge(
                    amenities[amenities[cluster_name]>-1].groupby(cluster_name).apply(lambda x: np.average(x['x'], weights=x['weight'])).reset_index().round(3).rename(columns={0:'x'}))
    result = result.merge(
                    amenities[amenities[cluster_name]>-1].groupby(cluster_name).apply(lambda x: np.average(x['y'], weights=x['weight'])).reset_index().round(3).rename(columns={0:'y'}))

    
    result['weight%'] = (result['weight'] / result['weight'].sum() * 100).round(1) 
    
    result = result.sort_values('weight', ascending=False).head(cantidad_clusters).reset_index().rename(columns={'index':'cluster'})
    
    result = gpd.GeoDataFrame(
                    result, geometry=gpd.points_from_xy(result['x'], result['y']), crs=city_crs).to_crs(4326)
    
    
    if show_map:
        fig = plt.figure(figsize=(15,15), dpi=100)

        ax = fig.add_subplot(2,2,1)
        amenities.to_crs(3857).plot(ax=ax, alpha=0)
        amenities[amenities[cluster_name]>-1].to_crs(3857).plot(ax=ax, column=cluster_name, categorical = True, lw=.1, alpha=.4)
        ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=3)
        ax.set_title(cluster_name, fontsize=8)
        ax.axis('off');


        ax = fig.add_subplot(2,2,2)
        amenities.to_crs(3857).plot(ax=ax, alpha=0)
        result.to_crs(3857).plot(ax=ax, column=cluster_name, categorical = True, lw=.1, alpha=.6, markersize = result['weight%']*5)
        ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=3)
        ax.set_title(cluster_name, fontsize=8)
        ax.axis('off');
        

    return result[['cluster', 'weight', 'weight%', 'geometry']], scores, amenities
    


    
    
    
def create_matrix(origin,                   
                  destination, 
                  id_origin = '', 
                  id_destination = '', 
                  latlon = True,                  
                  normalize = False, 
                  res = 8,
                  duplicates = False):
    
    '''
    En base a un GeoDataFrame de origines y uno de destinos crea una matriz de origin-destino para el cálculo de las distancias y/o tiempos.
    Parámetros:
    
    origin = GeoDataFrame de origines    
    destination = GeoDataFrame de destinos
    id_origin = ID del GeoDataFrame de origenes (opcional)
    id_destination = ID del GeoDataFrame de destinos (opcional)
    latlon = Especifica si la función devuelve una matriz con campos de latitud o longitud o con un campo de origin y uno de destino con lat/lon tipo texto (ej. "-31.06513, -64.28675")
    normalize = Si normalize es True, toma los origines y destinos desde el centroide del hexágono h3 con la resolución correspondiente. Este permite minimizar la cantidad de 
    consultas de tiempos de viaje en observaciones que están a una distancia muy cercana. Esto permite reducir los costos de consulta en Google Maps.
    res = Resolución de los hexágonos h3 si normalize = True.
    duplicates = False: Genera una matriz de origin y destinos sin duplicados. True: mantiene las tablas origin y destination sin modificar y puede traer duplicados si existen
    
    Salida: Matriz de origines y destinos
   '''
    
    origin_columns = origin.columns.tolist()
    destination_columns = destination.columns.tolist()
    origin_columns.remove('geometry')
    destination_columns.remove('geometry')
    
    lat_o, lon_o, lat_d, lon_d = 'lat_o', 'lon_o', 'lat_d', 'lon_d'
    
    origin = suph3.create_latlon(df = origin, normalize=True, res=res, hex='hex_o', var_latlon='origin', lat=lat_o, lon=lon_o).drop(['geometry'], axis=1)
    destination = suph3.create_latlon(df = destination, normalize=True, res=res, hex='hex_d', var_latlon='destination', lat=lat_d, lon=lon_d).drop(['geometry'], axis=1)
    
    cols_latlon_o, cols_latlon_d = [], []
    if latlon:        
        if not normalize:
            lat_o, lon_o, lat_d, lon_d = 'lat_o', 'lon_o', 'lat_d', 'lon_d'
        else:
            lat_o, lon_o, lat_d, lon_d = 'lat_o_norm', 'lon_o_norm', 'lat_d_norm', 'lon_d_norm'
        cols_latlon_o = [lat_o, lon_o]
        cols_latlon_d = [lat_d, lon_d]
    
    if not duplicates:
        if not normalize:
            origin = origin[['hex_o', 'origin', 'origin_norm']+cols_latlon_o].drop_duplicates().reset_index(drop=True)
            destination = destination[['hex_d', 'destination', 'destination_norm']+cols_latlon_d].drop_duplicates().reset_index(drop=True) 
        else:
            origin = origin[['hex_o', 'origin_norm']+cols_latlon_o].drop_duplicates().reset_index(drop=True)
            destination = destination[['hex_d', 'destination_norm']+cols_latlon_d].drop_duplicates().reset_index(drop=True) 
            
        cols = ['hex_o', 'hex_d']
    
    else:
        cols = []
        if len(id_origin)>0:
            cols += [id_origin]
        if len(id_destination)>0:
            cols += [id_destination]

        if id_origin != 'hex_o':
            cols += ['hex_o']
        if id_destination != 'hex_d':
            cols += ['hex_d']
                
    origin['aux'] = 1
    destination['aux'] = 1
    
    od_matrix = pd.merge(
                        origin,
                        destination,
                        on='aux').drop(['aux'], axis=1)
    
    del origin['aux']
    del destination['aux']
    
    if (not latlon)&(duplicates):
        od_matrix = od_matrix.drop(['lat_o', 'lon_o', 'lat_d', 'lon_d', 'lat_o_norm', 'lon_o_norm', 'lat_d_norm', 'lon_d_norm'], axis=1)
    
    for i in origin_columns:
        if (not i in cols)&(i in od_matrix.columns):
            cols += [i]
    
    for i in destination_columns:
        if (not i in cols)&(i in od_matrix.columns):
            cols += [i]
    
    
    for i in od_matrix.columns:
        if not i in cols:
            cols += [i]
    
    od_matrix = od_matrix[cols]
    
    
    return od_matrix


def measure_distances(idmatrix, node_from, node_to, G, lenx):
    '''
    Función de apoyo de measure_distances_osm
    '''
    
    if idmatrix % 2000 == 0:
        print(f'{datetime.datetime.now().strftime("%H:%M:%S")} processing {int(idmatrix)} / {lenx}')
        
    try:
        ret = nx.shortest_path_length(G, node_from, node_to, weight='length')        
    except:
        ret = np.nan
    return ret


def measure_distances_osm(origin, 
                          id_origin, 
                          destination, 
                          id_destination, 
                          driving = True,
                          walking = True,
                          normalize=False, 
                          res = 8,                           
                          processing='pandana', 
                          equipement_bring_closest = False, 
                          equipment_closest_qty = 3, 
                          closest_distance = [800, 1500, 2000],
                          equipment_type = '',
                          trips_file = 'trips_file_tmp',
                          current_path=Path() ):
    
    '''
    Calcula distancias en Open Street Maps para los modos de transporte entre dos capas cartográficas de origenes y destinos.
    En el caso de que alguna capa sea de polígonos, toma en cuenta el centroide de cada polígono.
    Devuelve una od_matrix de origenes y destinos con las distancias calculadas para cada modo.
    
    Parámetros
    origin = GeoDataFrame que representa a los puntos de origen
    id_origin = ID del GeoDataFrame origin
    destination = GeoDataFrame que representa a los puntos de destino
    id_destination = ID del GeoDataFrame destino
    driving = True    
    walking = True     
    processing= Se especifica donde procesar la consulta, puede ser a través de osmnx o pandana. Por defecto usa 'pandana'. Ambas busquedas utilizan Open Street Maps.
    equipement_bring_closest = Si trae solo los destinos más cercanos. False trae todos los destinos, True. Ordena por distancia y trae los más cercanos según la variable equipment_closes_qty
    equipment_closest_qty = Cantidad de destinos a traer si equipement_bring_closest es True
    equipment_type = Si hay más de un tipo de equipamiento en la tabla, por ejemplo Escuelas y Hospitales en una variable tipo_establecimientos (ej. equipment_type = 'tipo_establecimientos'). Se puede enviar una lista de variables.
    trips_file = Nombre de archivo temporal donde se guardan las consultas. Por defecto 'trips_file_tmp'
    current_path = Directorio de trabajo. Por defecto: Path()
    
    Salida: Matriz de origenes y destino con los campos nuevos de distancia calculados en Open Street Maps para los modos solicitados. 
    
    '''

    add_file = ''
    if driving: add_file+='_drive'
    if walking: add_file+='_walk'
    if normalize: add_file+='_norm'
    
    create_result_dir(current_path=current_path)
    
    modes = []
    if driving:
        modes += ['drive']
    if walking:
        modes += ['walk']
        
    if not normalize:
        lat_o, lon_o, lat_d, lon_d = 'lat_o', 'lon_o', 'lat_d', 'lon_d'
        geo_origin, geo_destination = 'origin', 'destination'
    else:
        lat_o, lon_o, lat_d, lon_d = 'lat_o_norm', 'lon_o_norm', 'lat_d_norm', 'lon_d_norm'
        geo_origin, geo_destination = 'origin_norm', 'destination_norm'
    
    
    od_matrix_all = create_matrix(origin, 
                                  destination, 
                                  id_origin=id_origin, 
                                  id_destination = id_destination, 
                                  normalize = False, 
                                  res = res, 
                                  duplicates = True, 
                                  latlon = False)
    
    od_matrix = create_matrix(origin, 
                              destination, 
                              id_origin = id_origin, 
                              id_destination = id_destination, 
                              latlon = True, 
                              normalize = normalize, res=res)
    
    bounds_tmp = pd.concat([origin['geometry'], destination['geometry']])
    
    xmin, ymin, xmax, ymax = bounds_tmp.total_bounds
    
    trips = pd.DataFrame([])
    if Path(trips_file).is_file():                
        trips = pd.read_csv(trips_file)
        
        trips['osm'] = 1

        if (geo_origin in trips.columns)&(geo_destination in trips.columns):

            od_matrix = od_matrix.merge(trips[[geo_origin, 
                                               geo_destination, 
                                               'osm']],
                                                how='left', 
                                                on=[geo_origin, 
                                                    geo_destination])
            
            od_matrix = od_matrix[od_matrix.osm.isna()]
            del od_matrix['osm']
            del trips['osm']
        
        

    if len(od_matrix) > 0:
        var_distances = []

        for mode in modes:

            print(f'Coords OSM {mode} - Download map') 
            print('')

            if processing != 'pandana':

                G = ox.graph_from_bbox(ymax,
                                       ymin,
                                       xmax, 
                                       xmin, 
                                       network_type=mode)

                G = ox.add_edge_speeds(G)
                G = ox.add_edge_travel_times(G)

                nodes_from = ox.distance.nearest_nodes(G, 
                                                       od_matrix[lon_o].values, 
                                                       od_matrix[lat_o].values, 
                                                       return_dist=True)

                od_matrix['node_from'] = nodes_from[0]

                nodes_to = ox.distance.nearest_nodes(G, 
                                                     od_matrix[lon_d].values, 
                                                     od_matrix[lat_d].values, 
                                                     return_dist=True)

                od_matrix['node_to'] = nodes_to[0]

                od_matrix = od_matrix.reset_index().rename(columns={'index':'idmatrix'})
                od_matrix[f'distance_osm_{mode}'] = od_matrix.apply(lambda x : measure_distances(x['idmatrix'],
                                                                                                 x['node_from'], 
                                                                                                 x['node_to'], 
                                                                                                 G = G, 
                                                                                                 lenx = len(od_matrix)), 
                                                           axis=1)
            else:
                
                network = osm_pandana.pdna_network_from_bbox(ymin, xmin, ymax,  xmax, network_type=mode)  

                od_matrix['node_from'] = network.get_node_ids(od_matrix[lon_o], od_matrix[lat_o]).values
                od_matrix['node_to'] = network.get_node_ids(od_matrix[lon_d], od_matrix[lat_d]).values
                od_matrix[f'distance_osm_{mode}'] = network.shortest_path_lengths(od_matrix['node_to'].values, od_matrix['node_from'].values) 
            
            var_distances += [f'distance_osm_{mode}']
            od_matrix[f'distance_osm_{mode}'] = (od_matrix[f'distance_osm_{mode}'] / 1000).round(2)
            print('')

        od_matrix = od_matrix[['hex_o', 'hex_d', geo_origin, geo_destination]+var_distances]
        

    od_matrix = pd.concat([trips, od_matrix], ignore_index=True)
    
    if 'lat_o' in od_matrix.columns:
        od_matrix = od_matrix.drop(['lat_o', 'lon_o', 'lat_d', 'lon_d'], axis=1)
    if 'lat_o_norm' in od_matrix.columns:
        od_matrix = od_matrix.drop(['lat_o_norm', 'lon_o_norm', 'lat_d_norm', 'lon_d_norm'], axis=1)
    if (not normalize)&('origin_norm' in od_matrix.columns):
        od_matrix = od_matrix.drop(['origin_norm', 'destination_norm'], axis=1)

    
    od_matrix.to_csv(trips_file, index=False)

    od_matrix = od_matrix_all.merge(od_matrix, how='left', on=['hex_o', 'hex_d', geo_origin, geo_destination])
    
    if equipement_bring_closest:
        
        if len(equipment_type) == 0:
            od_matrix['equipment_type_tmp'] = 1
            equipment_type = ['equipment_type_tmp']

        if type(equipment_type) == str:
            equipment_type = [equipment_type]

        od_matrix = od_matrix.sort_values(equipment_type + [id_origin, 'distance_osm_walk'])
        od_matrix['distance_order'] = od_matrix.groupby(equipment_type + [id_origin]).transform('cumcount')
        od_matrix['distance_order'] = od_matrix['distance_order'] + 1
        
        
        if len(closest_distance)>0:
            if type(closest_distance) == str:
                closest_distance = [closest_distance]
            for i in closest_distance:
                od_matrix_tmp = od_matrix[od_matrix.distance_osm_walk<=(i/1000)].groupby(equipment_type + [id_origin], as_index=False).size().rename(columns={'size':f'qty_est_{i}m'})                
                od_matrix = od_matrix.merge(od_matrix_tmp, on=equipment_type + [id_origin], how='left')
                od_matrix[f'qty_est_{i}m'] = od_matrix[f'qty_est_{i}m'].fillna(0).astype(int)
                
        od_matrix = od_matrix[od_matrix['distance_order']<=equipment_closest_qty].reset_index(drop=True)

        if 'equipment_type_tmp' in od_matrix.columns:
            od_matrix = od_matrix.drop(['equipment_type_tmp'], axis=1)
    
    
    print('Proceso finalizado')
    return od_matrix

def bring_osm(gdf, tags = {'leisure':True, 'aeroway':True, 'natural':True}, list_types = '' ):
    '''
    Trae elementos de Open Street Maps usando la librería osmnx
    
    Parámetros:
    gdf = Capa de referencia para delimitar el área de elementos a traer
    tags = Elementos del mapa a traer. Una lista completa de elementos en Open Street Maps se puede consultar en: https://wiki.openstreetmap.org/wiki/Map_features
    list_types: Se descargan puntos, líneas, polígonos y multipolígonos. list_types define el tipo de elementos a traer.
    
    Salida: GeoDataFrame con objetos espaciales de Open Street Maps
    '''

    xmin,ymin,xmax,ymax = gdf.total_bounds
    osm = ox.geometries_from_bbox(ymax,ymin,xmax, xmin, tags).reset_index()

    for i in tags:    
        osm.loc[osm[i].notna(), 'tag'] = i
        osm.loc[osm[i].notna(), 'tag_type'] = osm.loc[osm[i].notna(), i]

    osm = osm[['osmid', 'tag', 'tag_type', 'name', 'geometry']]
    if len(list_types) > 0:
        osm = osm.loc[(osm.geometry.type.isin(list_types)), :] 

    return osm
    
    

    
    
    
def trip_info_googlemaps(coord_origin, 
                         coord_destin, 
                         trip_datetime, 
                         trip_datetime_corrected,
                         gmaps,
                         transit = True, 
                         driving = True, 
                         walking = False, 
                         bicycling = False):
    
    '''
    Realiza una consulta a la API de Google Maps teniendo en cuenta un origen y un destino
    
    Parámetros:
    coord_origin = latitud/longitud (xy) del origen del viaje en formato texto. Ej. '-31.4247, -64.15922'
    coord_destin = = latitud/longitud (xy) del destino del viaje en formato texto. Ej. '-31.4247, -64.15922'
    trip_datetime = Fecha y hora de la consulta del viaje en formato datetime
    gmaps = Objeto gmaps para acceder a la API de googlemaps
    transit = Para obtener distancias y tiempos en transporte público (True/False). Por defecto True.
    driving = Para obtener distancias y tiempos en automovil (True/False). Por defecto True.
    walking = Para obtener distancias y tiempos caminando (True/False). Por defecto False.
    bicycling = Para obtener distancias y tiempos en bicicleta (True/False). Por defecto False.
    
    Salida: nuevo dataframe con las distancias y tiempos calculados para ese origen y destino
    
    '''
    
    alternatives = False
    
    
    
    Transit, Driving, Walking, Bicycling = '', '', '', ''
    
    if transit:
        Transit = gmaps.directions(coord_origin,
                                   coord_destin,
                                   mode='transit', 
                                   departure_time=trip_datetime_corrected, 
                                   alternatives=alternatives, 
                                   traffic_model="best_guess")
        

    if driving:
        Driving = gmaps.directions(coord_origin,
                                   coord_destin,
                                   mode='driving', 
                                   departure_time=trip_datetime_corrected, 
                                   alternatives=alternatives, 
                                   traffic_model="best_guess")

    if walking:
        Walking = gmaps.directions(coord_origin,
                                   coord_destin,
                                   mode='walking', 
                                   departure_time=trip_datetime_corrected, 
                                   alternatives=alternatives, 
                                   traffic_model="best_guess")

    
    if bicycling:
        Bicycling = gmaps.directions(coord_origin,
                                   coord_destin,
                                   mode='bicycling', 
                                   departure_time=trip_datetime_corrected, 
                                   alternatives=alternatives, 
                                   traffic_model="best_guess")


    transit_trips, driving_trips, walking_trips, bicycling_trips = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])

    # Transit
    v1, v2, v3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])

    if len(Transit) != 0:
        t1, t2, t3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
        t1 = pd.DataFrame.from_dict(pd.json_normalize(Transit))
        for i in range(0, len(t1)):

            t1.loc[i,"trip_datetime"] = trip_datetime                
            t1.loc[i,"geo_origin"] = coord_origin
            t1.loc[i,"geo_destination"] = coord_destin
            t1.loc[i,"alternative"] = i

            t2 = pd.DataFrame.from_dict(pd.json_normalize(Transit[i]['legs']))
            for l in range(0, len(t2)):
                t2.loc[l,"trip_datetime"] = trip_datetime                    
                t2.loc[l,"geo_origin"] = coord_origin
                t2.loc[l,"geo_destination"] = coord_destin
                t2.loc[l,"alternative"] = i
                t2.loc[l,"legs"] = l

                t3 = pd.DataFrame.from_dict(pd.json_normalize(Transit[i]['legs'][l]['steps']))        
                for s in range(0, len(t3)):
                    t3.loc[s,"trip_datetime"] = trip_datetime                        
                    t3.loc[s,"geo_origin"] = coord_origin
                    t3.loc[s,"geo_destination"] = coord_destin
                    t3.loc[s,"alternative"] = i
                    t3.loc[s,"legs"] = l
                    t3.loc[s,"steps"] = s

                v3 = pd.concat([v3, t3], ignore_index=True)
            v2 = pd.concat([v2, t2], ignore_index=True)
        v1 = pd.concat([v1, t1], ignore_index=True)
        
        
        
        if 'transit_details.line.vehicle.name' not in v3.columns:
            v3['transit_details.line.vehicle.name'] = np.nan
            
        if 'transit_details.line.vehicle.type' not in v3.columns:
            v3['transit_details.line.vehicle.type'] = np.nan
            


        v3 = v3[['trip_datetime', 
            'geo_origin', 
            'geo_destination', 
            'alternative', 
            'steps', 
            'travel_mode', 
            'distance.value', 
            'duration.value',              
            'transit_details.line.vehicle.name',
            'transit_details.line.vehicle.type'
               ]]

        v3 = pd.concat([v3, 
                        pd.get_dummies(v3['travel_mode']),
                        pd.get_dummies(v3['transit_details.line.vehicle.type'])], axis=1)
        
        transit_modes = []
        for i in v3['travel_mode'].unique().tolist():
            v3[f'transit_{i.lower()}.distance'] = v3['distance.value'] * v3[i]
            v3[f'transit_{i.lower()}.duration'] = v3['duration.value'] * v3[i]
            v3[f'transit_{i.lower()}.steps'] = 1 * v3[i]

            transit_modes += [f'transit_{i.lower()}.distance', f'transit_{i.lower()}.duration', f'transit_{i.lower()}.steps']

        steps = v3.groupby(['trip_datetime', 'geo_origin', 'geo_destination', 'alternative'])[transit_modes].sum().reset_index()
        
        if 'departure_time.text' not in v2.columns:
            v2['departure_time.text'] = np.nan
            
        if 'arrival_time.text' not in v2.columns:
            v2['arrival_time.text'] = np.nan
        

        transit_trips = v2[['trip_datetime', 
                            'geo_origin', 
                            'geo_destination', 
                            'alternative',                             
                            'departure_time.text',                        
                            'arrival_time.text',
                            'distance.value',                    
                            'duration.value',                     
                            ]].copy()

        transit_trips = transit_trips.merge(steps, how='left', on=['trip_datetime', 
                                                                'geo_origin', 
                                                                'geo_destination', 
                                                                'alternative'])
        transit_trips = transit_trips.merge(
                                v3.loc[(v3.steps==0)&(v3.travel_mode=='WALKING'), ['trip_datetime', 
                                                                                   'geo_origin', 
                                                                                   'geo_destination', 
                                                                                   'alternative', 
                                                                                   'distance.value', 
                                                                                   'duration.value']].rename(columns={'distance.value': 'transit_walking.distance.origin',
                                                                                                                      'duration.value': 'transit_walking.duration.origin'}),
                                how='left', 
                                on=['trip_datetime', 'geo_origin', 'geo_destination', 'alternative'])

        transit_trips.columns = [i.replace('.value', '') for i in transit_trips.columns ]
        transit_trips.columns = [i.replace('.text', '') for i in transit_trips.columns ]

        transit_trips = transit_trips.rename(columns={'arrival_time': 'transit_arrival_time', 
                                                     'departure_time':'transit_departure_time',
                                                     'distance':'transit_distance',
                                                     'duration': 'transit_duration'})


    # Driving
    d1, d2, d3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
    if len(Driving) != 0:    
        t1, t2, t3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
        t1 = pd.DataFrame.from_dict(pd.json_normalize(Driving))
        for i in range(0, len(t1)):
            t1.loc[i,"trip_datetime"] = trip_datetime                
            t1.loc[i,"geo_origin"] = coord_origin
            t1.loc[i,"geo_destination"] = coord_destin
            t1.loc[i,"alternative"] = i

            t2 = pd.DataFrame.from_dict(pd.json_normalize(Driving[i]['legs']))
            for l in range(0, len(t2)):
                t2.loc[l,"trip_datetime"] = trip_datetime                
                t2.loc[l,"geo_origin"] = coord_origin
                t2.loc[l,"geo_destination"] = coord_destin
                t2.loc[l,"alternative"] = i
                t2.loc[l,"legs"] = l

                t3 = pd.DataFrame.from_dict(pd.json_normalize(Driving[i]['legs'][l]['steps']))        
                for s in range(0, len(t3)):
                    t3.loc[s,"trip_datetime"] = trip_datetime                
                    t3.loc[s,"geo_origin"] = coord_origin
                    t3.loc[s,"geo_destination"] = coord_destin
                    t3.loc[s,"alternative"] = i
                    t3.loc[s,"legs"] = l
                    t3.loc[s,"steps"] = s

                d3 = pd.concat([d3, t3], ignore_index=True)

            d2 = pd.concat([d2, t2], ignore_index=True)
        d1 = pd.concat([d1, t1], ignore_index=True)
        
        if 'duration_in_traffic.value' not in d2.columns:            
            d2['duration_in_traffic.value'] = np.nan

        
        driving_trips = d2[['trip_datetime',
                           'geo_origin', 'geo_destination', 'alternative',                   
                           'distance.value', 
                           'duration.value', 
                           'duration_in_traffic.value']]
        driving_trips.columns = [i.replace('.value', '') for i in driving_trips.columns ]
        driving_trips = driving_trips.rename(columns={'distance':'driving_distance', 'duration': 'driving_duration', 'duration_in_traffic':'driving_duration_in_traffic'})

    # Walking
    w1, w2, w3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
    if len(Walking) != 0:
        t1, t2, t3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
        t1 = pd.DataFrame.from_dict(pd.json_normalize(Walking))
        for i in range(0, len(t1)):
            t1.loc[i,"trip_datetime"] = trip_datetime                
            t1.loc[i,"geo_origin"] = coord_origin
            t1.loc[i,"geo_destination"] = coord_destin
            t1.loc[i,"alternative"] = i

            t2 = pd.DataFrame.from_dict(pd.json_normalize(Walking[i]['legs']))
            for l in range(0, len(t2)):
                t2.loc[l,"trip_datetime"] = trip_datetime                
                t2.loc[l,"geo_origin"] = coord_origin
                t2.loc[l,"geo_destination"] = coord_destin
                t2.loc[l,"alternative"] = i
                t2.loc[l,"legs"] = l

                t3 = pd.DataFrame.from_dict(pd.json_normalize(Walking[i]['legs'][l]['steps']))        
                for s in range(0, len(t3)):
                    t3.loc[s,"trip_datetime"] = trip_datetime                
                    t3.loc[s,"geo_origin"] = coord_origin
                    t3.loc[s,"geo_destination"] = coord_destin
                    t3.loc[s,"alternative"] = i
                    t3.loc[s,"legs"] = l
                    t3.loc[s,"steps"] = s

                w3 = pd.concat([w3, t3], ignore_index=True)
            w2 = pd.concat([w2, t2], ignore_index=True)
        w1 = pd.concat([w1, t1], ignore_index=True)

        walking_trips = w2[['trip_datetime',
                            'geo_origin', 
                            'geo_destination', 
                            'alternative', 
                            'distance.value',     
                            'duration.value']]
        walking_trips.columns = [i.replace('.value', '') for i in walking_trips.columns ]
        walking_trips = walking_trips.rename(columns={'distance':'walking_distance', 
                                                      'duration': 'walking_duration'})

    # Biciclying
    b1, b2, b3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])

    if len(Bicycling) != 0:
        t1, t2, t3 = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
        t1 = pd.DataFrame.from_dict(pd.json_normalize(Bicycling))
        for i in range(0, len(t1)):
            t1.loc[i,"trip_datetime"] = trip_datetime                
            t1.loc[i,"geo_origin"] = coord_origin
            t1.loc[i,"geo_destination"] = coord_destin
            t1.loc[i,"alternative"] = i

            t2 = pd.DataFrame.from_dict(pd.json_normalize(Bicycling[i]['legs']))
            for l in range(0, len(t2)):
                t2.loc[l,"trip_datetime"] = trip_datetime                
                t2.loc[l,"geo_origin"] = coord_origin
                t2.loc[l,"geo_destination"] = coord_destin
                t2.loc[l,"alternative"] = i
                t2.loc[l,"legs"] = l

                t3 = pd.DataFrame.from_dict(pd.json_normalize(Bicycling[i]['legs'][l]['steps']))        
                for s in range(0, len(t3)):
                    t3.loc[s,"trip_datetime"] = trip_datetime                
                    t3.loc[s,"geo_origin"] = coord_origin
                    t3.loc[s,"geo_destination"] = coord_destin
                    t3.loc[s,"alternative"] = i
                    t3.loc[s,"legs"] = l
                    t3.loc[s,"steps"] = s

                b3 = pd.concat([b3, t3], ignore_index=True)
            b2 = pd.concat([b2, t2], ignore_index=True)
        b1 = pd.concat([b1, t1], ignore_index=True)

        bicycling_trips = b2[['trip_datetime',
                            'geo_origin', 
                            'geo_destination', 'alternative',
                            'distance.value', 
                            'duration.value']]
        bicycling_trips.columns = [i.replace('.value', '') for i in bicycling_trips.columns ]
        bicycling_trips = bicycling_trips.rename(columns={'distance':'bicycling_trips_distance', 
                                                      'duration': 'bicycling_trips_duration'})

    if len(transit_trips)==0:    
        transit_trips['trip_datetime'] = np.nan
        transit_trips['geo_origin'] = np.nan
        transit_trips['geo_destination'] = np.nan
        transit_trips['alternative']  = np.nan
    if len(driving_trips)==0:
        driving_trips['trip_datetime'] = np.nan
        driving_trips['geo_origin'] = np.nan
        driving_trips['geo_destination'] = np.nan
        driving_trips['alternative']  = np.nan
    if len(walking_trips)==0:
        walking_trips['trip_datetime'] = np.nan
        walking_trips['geo_origin'] = np.nan
        walking_trips['geo_destination'] = np.nan
        walking_trips['alternative']  = np.nan
    if len(bicycling_trips)==0:
        bicycling_trips['trip_datetime'] = np.nan
        bicycling_trips['geo_origin'] = np.nan
        bicycling_trips['geo_destination'] = np.nan
        bicycling_trips['alternative']  = np.nan


    if len(transit_trips) == 0:
        transit_trips = driving_trips.copy()
    else:
        transit_trips = transit_trips.merge(driving_trips, how='left', on=[ 'trip_datetime',
                                                                            'geo_origin', 
                                                                            'geo_destination', 
                                                                            'alternative' ])
    if len(transit_trips) == 0:
        transit_trips = walking_trips.copy()
    else:
        transit_trips = transit_trips.merge(walking_trips, how='left', on=[ 'trip_datetime',
                                                                            'geo_origin', 
                                                                            'geo_destination', 
                                                                            'alternative' ])
    if len(transit_trips) == 0:
        transit_trips = bicycling_trips.copy()
    else:
        transit_trips = transit_trips.merge(bicycling_trips, how='left', on=[ 'trip_datetime',
                                                                            'geo_origin', 
                                                                            'geo_destination', 
                                                                            'alternative' ])

    for i in transit_trips.columns:
        if i.find('duration') > -1:
            transit_trips[i] = (transit_trips[i] / 60).round(2)
        if i.find('distance') > -1:
            transit_trips[i] = (transit_trips[i] / 1000).round(2)
        if i.find('.') > -1:
            transit_trips = transit_trips.rename(columns={i: i.replace('.', '_')})
            
    if not alternatives:        
        transit_trips = transit_trips.drop(['alternative'], axis=1)
        
    if len(transit_trips) == 0:
        transit_trips = pd.DataFrame([[trip_datetime, coord_origin, coord_destin]], columns=['trip_datetime', 'origin_norm', 'destination_norm'])
        
    return transit_trips

def gmaps_matrix(od_matrix, geo_origin, geo_destination, trip_datetime, mode, gmaps):
    '''
    Realiza una consulta a la API de Google Maps teniendo en cuenta un dataframe de origenes y destinos. 
    Las consultas se realizan para una matriz de origenes y destinos y devuelve solo tiempos y distancias del viaje.
    
    Parámetros:
    od_matrix = Matriz de origenes y destinos. Debe contener un campo geo_origen y otro geo_destination en formato lat/lon texto para realizar el proceso
    trip_datetime = Fecha y hora de la consulta del viaje en formato datetime
    mode = Según el modo que se quiera consultar, las opciones son: transit, driving, walking, bicycling
    gmaps = Objeto gmaps para acceder a la API de googlemaps
    
    Salida: Matriz de origenes y destinos con los resultados de la consulta a Google Maps
    '''

    trip_datetime_corrected = sup.correct_datetime(trip_datetime, lat=float(od_matrix[geo_origin].head(1).values[0].split(',')[0]), lon=float(od_matrix[geo_origin].head(1).values[0].split(',')[1]))
    
    list_origin = od_matrix[geo_origin].unique().tolist()
    list_destination = od_matrix[geo_destination].unique().tolist()

    try:
        result = gmaps.distance_matrix(list_origin,
                                       list_destination, 
                                       mode=mode,
                                       departure_time=trip_datetime_corrected)
    
    
        t1 = pd.DataFrame.from_dict(pd.json_normalize(result))
        origin_addresses = t1['origin_addresses'][0]
        destination_addresses = t1['destination_addresses'][0]

        trips = pd.DataFrame([])
        n = 0
        for x in t1['rows'][0]:     
            tx = pd.DataFrame.from_dict(pd.json_normalize(x['elements']))    

            tx['origin_address'] = origin_addresses[n]
            tx[geo_origin] = list_origin[n]
            tx = pd.concat([tx,
                   pd.DataFrame(list_destination, columns=[geo_destination]),
                   pd.DataFrame(destination_addresses, columns=['destination_address'])], axis=1)

            tx['trip_datetime'] = trip_datetime

            if not 'distance.value' in tx.columns:
                tx['distance.value'] = np.nan
            if not 'duration.value' in tx.columns:
                tx['duration.value'] = np.nan

            trips = trips.append(tx, ignore_index=True)
            n+=1

        if len(trips)>0:
            var_list = ['trip_datetime', geo_origin, geo_destination, 'origin_address', 'destination_address', 'distance.value', 'duration.value']
            var_list_new = ['trip_datetime', geo_origin, geo_destination, 'origin_address', 'destination_address', f'{mode}_distance', f'{mode}_duration']
            if mode == 'driving':
                trips['duration_in_traffic.value'] = round(trips['duration_in_traffic.value'] / 60, 2)
                var_list = var_list+['duration_in_traffic.value']
                var_list_new = var_list_new + [f'{mode}_duration_in_traffic']

            trips = trips[var_list]
            trips.columns = var_list_new

            trips[ f'{mode}_distance'] = round(trips[ f'{mode}_distance'] / 1000,2)
            trips[ f'{mode}_duration'] = round(trips[ f'{mode}_duration'] / 60,2)

    except:
        print('')
        print(f'** Error en la API de Google Maps. Revise que esté funcionando correctamente la key utilizada _ modo {mode}')
        print('')
        trips = pd.DataFrame([])
    
    return trips 
    
    
def save_key(key, Qty, current_path=Path(), key_file = 'save_key.csv'):
    '''
    Guarda cantidad de consultas en Google Maps en un archivo auxiliar    
    '''
    save_key_file = current_path / 'tmp' / key_file
    n = 0
    new_key = ''
    for i in range(0,39, 5):
        if not n % 2:
            new_key += key[i:i+5]
        else:
            new_key += '....'    
        n += 1
    new_key

    try:

        tmp = current_path / 'tmp'
        if save_key_file.is_file():        
            save_key = pd.read_csv(save_key_file)  
        else:
            save_key = pd.DataFrame([[new_key, str(date.today())[:7], 0]], columns=['key', 'month', 'Qty'])

        if len(save_key[(save_key.key==new_key)&(save_key.month==str(date.today())[:7])]) == 0:
            save_key = save_key.append(pd.DataFrame([[new_key, str(date.today())[:7], 0]], columns=['key', 'month', 'Qty']), ignore_index=True)

        save_key.loc[(save_key.key==new_key)&(save_key.month==str(date.today())[:7]), 'Qty'] += Qty

        save_key['usd'] = save_key['Qty'] * .01

        save_key.to_csv(save_key_file, index=False)
    except:
        print('except')
        save_key = pd.DataFrame([[new_key, str(date.today())[:7], 0]], columns=['key', 'month', 'Qty'])
        save_key.to_csv(save_key_file, index=False)
    
    save_key['X'] = ''    
    save_key.loc[(save_key.key==new_key)&(save_key.month==str(date.today())[:7]), 'X'] = 'X'
    
    return save_key
    
    
    
def trip_info_googlemaps_matrix(od_matrix, 
                                geo_origin,
                                geo_destination,
                                trip_datetime, 
                                gmaps,
                                transit = True, 
                                driving = True, 
                                walking = False, 
                                bicycling = False):
    '''
    Función auxiliar para la creación de la matrix de origenes y destinos con distintos modos de transporte    
    '''
    
    transit_trips, driving_trips, walking_trips, bicycling_trips = pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([]), pd.DataFrame([])
    
    if transit:
        transit_trips = gmaps_matrix(od_matrix, geo_origin, geo_destination, trip_datetime, 'transit', gmaps)
    if driving:
        driving_trips = gmaps_matrix(od_matrix, geo_origin, geo_destination, trip_datetime, 'driving', gmaps)
    if walking:
        walking_trips = gmaps_matrix(od_matrix, geo_origin, geo_destination, trip_datetime, 'walking', gmaps)
    if bicycling:
        bicycling_trips = gmaps_matrix(od_matrix, geo_origin, geo_destination, trip_datetime, 'bicycling', gmaps)

    if len(transit_trips)==0:    
        transit_trips['trip_datetime'] = np.nan
        transit_trips[geo_origin] = np.nan
        transit_trips[geo_destination] = np.nan
        transit_trips['origin_address']  = np.nan
        transit_trips['destination_address']  = np.nan
    if len(driving_trips)==0:
        driving_trips['trip_datetime'] = np.nan
        driving_trips[geo_origin] = np.nan
        driving_trips[geo_destination] = np.nan
        driving_trips['origin_address']  = np.nan
        driving_trips['destination_address']  = np.nan
    if len(walking_trips)==0:
        walking_trips['trip_datetime'] = np.nan
        walking_trips[geo_origin] = np.nan
        walking_trips[geo_destination] = np.nan
        walking_trips['origin_address']  = np.nan
        walking_trips['destination_address']  = np.nan
    if len(bicycling_trips)==0:
        bicycling_trips['trip_datetime'] = np.nan
        bicycling_trips[geo_origin] = np.nan
        bicycling_trips[geo_destination] = np.nan
        bicycling_trips['origin_address']  = np.nan
        bicycling_trips['destination_address']  = np.nan


    if len(transit_trips) == 0:
        transit_trips = driving_trips.copy()
    else:    
        driving_trips = driving_trips.drop(['origin_address', 'destination_address'], axis=1)
        transit_trips = transit_trips.merge(driving_trips, how='left', on=[ 'trip_datetime',
                                                                            geo_origin, 
                                                                            geo_destination])
    if len(transit_trips) == 0:
        transit_trips = walking_trips.copy()
    else:
        walking_trips = walking_trips.drop(['origin_address', 'destination_address'], axis=1)
        transit_trips = transit_trips.merge(walking_trips, how='left', on=[ 'trip_datetime',
                                                                            geo_origin, 
                                                                            geo_destination])
    if len(transit_trips) == 0:
        transit_trips = bicycling_trips.copy()
    else:    
        bicycling_trips = bicycling_trips.drop(['origin_address', 'destination_address'], axis=1)
        transit_trips = transit_trips.merge(bicycling_trips, how='left', on=[ 'trip_datetime',
                                                                            geo_origin, 
                                                                            geo_destination])
    return transit_trips




def trips_gmaps_process(od_matrix, 
                        geo_origin,
                        geo_destination,
                        trip_datetime,                 
                        key, 
                        transit=True,
                        driving=True,
                        walking=False,
                        bicycling=False,
                        normalize=True,
                        res = 8,
                        trips_file = 'trips_file_tmp',
                        full_day=False,
                        only_distance_duration=False,
                        current_path=Path()):
    
    '''
    Función auxiliar para el proceso de consulta a la API de Google Maps. Se llama de trips_gmaps_from_od y de trips_gmaps_from_matrix
    '''
    
# if True:
#     od_matrix = od_matrix_tmp.head(2).copy() 
#     geo_origin=geo_origin
#     geo_destination=geo_destination 
#     trip_datetime = list_trip_datetime 
#     transit=transit
#     driving=driving
#     walking=walking
#     bicycling=bicycling
#     normalize=normalize
#     res = res    
#     full_day=full_day
#     only_distance_duration=only_distance_duration
#     current_path=current_path
#     trips_file = 'trips_file_tmp'

    create_result_dir(current_path=current_path)    
    
    modos = ''
    answer = ''
    if transit: modos += 'Transporte Público, '
    if driving: modos += 'Automovil, '
    if walking: modos += 'Caminata, '
    if bicycling: modos += 'Bicicleta, '
    modos = modos[:-2]

    if transit: trips_file += '_transit'
    if driving: trips_file += '_drive'
    if walking: trips_file += '_walk'
    if bicycling: trips_file += '_bike'

    if only_distance_duration:
        trips_file += '_matrix'
    if normalize:
        trips_file += '_norm'
        
    if 'trip_datetime' in od_matrix.columns: del od_matrix['trip_datetime']
        
    list_trip_datetime = trip_datetime
    if type(trip_datetime) != list:
        list_trip_datetime = [list_trip_datetime]
        
    od_matrix_ = od_matrix.copy()
    od_matrix = pd.DataFrame([])
    for i in list_trip_datetime:
        od_matrix_['trip_datetime'] = pd.to_datetime(i)
        od_matrix = pd.concat([od_matrix, od_matrix_], ignore_index=True)    
    
    if full_day:
        
        od_matrix_full_day = pd.DataFrame([])
        for n in range(0, 24):                    
            od_matrix['trip_datetime'] = od_matrix.trip_datetime.apply(lambda dt: dt.replace(hour=n))
            od_matrix['hour'] = n
            od_matrix_full_day = pd.concat([od_matrix_full_day,
                                            od_matrix], ignore_index=True)
        od_matrix_all = od_matrix_full_day.drop_duplicates().copy()
        
    else:
        od_matrix_all = od_matrix.copy()
    od_matrix_all = od_matrix_all.drop_duplicates()
    
    trips_all = pd.DataFrame([])
    for trip_datetime in list_trip_datetime:
        trips_file_ = current_path / 'tmp' / Path(f'{trips_file}_{str(trip_datetime)[:10]}.csv'.replace(':', '_').replace(' ', '_'))

        if Path(trips_file_).is_file():             
            
            trips_ = pd.read_csv(trips_file_) 
            trips_['trip_datetime'] = pd.to_datetime(trips_['trip_datetime'])
            
            trips_['gmaps'] = 1
            trips_all = pd.concat([trips_all, trips_], ignore_index=True)
            
    if len(trips_all)>0:

        od_matrix_agg = od_matrix_all.merge(trips_all[[geo_origin, 
                                                       geo_destination, 
                                                       'trip_datetime', 
                                                       'gmaps']],
                                            how='left', 
                                            on=[geo_origin, 
                                                geo_destination, 
                                                'trip_datetime'])

        
        od_matrix_agg = od_matrix_agg[(od_matrix_agg.gmaps.isna())&(od_matrix_agg[geo_origin]!=od_matrix_agg[geo_destination])]
    
        del od_matrix_agg['gmaps']
        del trips_all['gmaps']
    else:
        od_matrix_agg = od_matrix_all.copy()

    trips_all_new = trips_all.copy()
    
    if len(od_matrix_agg) > 0:
        if 'hour' in od_matrix_agg:
            od_matrix_agg = od_matrix_agg.sort_values(['trip_datetime', 'hour'])
            del od_matrix_agg['hour']
        else:
            od_matrix_agg = od_matrix_agg.sort_values(['trip_datetime'])
                
        try:
            gmaps = googlemaps.Client(key)
            gmaps_ok = True

        except:
            print('')
            print('Hay un error en la configuración de la API de Google Maps')
            print('Verifique que el API Key sea correcto y esté valido y que haya conexión a internet')
            print('')
            gmaps_ok = False
            od_matrix = pd.DataFrame([])

        if gmaps_ok:

            if trip_datetime > datetime.datetime.now():

                weekDaysMapping = ("Lunes", 
                                   "Martes",
                                   "Miércoles", 
                                   "Jueves",
                                   "Viernes", 
                                   "Sábado",
                                   "Domingo")
                monthsMapping = ("", 
                                 "Enero", 
                                 "Febrero", 
                                 "Marzo", 
                                 "Abril", 
                                 "Mayo", 
                                 "Junio", 
                                 "Julio", 
                                 "Agosto", 
                                 "Septiembre", 
                                 "Octubre", 
                                 "Noviembre", 
                                 "Diciembre")

                cost = transit * .01 + driving * .015 + walking * .01 + bicycling*.01
                qty_queries = len(od_matrix_agg[(od_matrix_agg[geo_origin]!=od_matrix_agg[geo_destination])])*(transit * 1 + driving * 1 + walking * 1 + bicycling*1)
                qty_queries_str = '{:,}'.format(qty_queries)
                len_matrix = len(od_matrix)
                len_matrix = '{:,}'.format(len_matrix)
                
                print(f' Para una matriz de origenes y destinos de {len_matrix} viajes se van a realizar {qty_queries_str} consultas en la Api de Google Maps')
                print('')
                print(f' Se van a consultar los modos: {modos} a un costo estimado de USD {round(len(od_matrix_agg)* cost, 2)}')                    
                print('')
                
                print('Las consultas se realizarán para los siguientes días:')
                if not full_day:
                    for ii in od_matrix_agg.trip_datetime.unique():                    
                        ii = pd.to_datetime(ii)
                        print(f'        {weekDaysMapping[ii.weekday()]} {ii.day} de {monthsMapping[ii.month]} de {ii.year} a las {str(ii.hour).zfill(2)}:{str(ii.minute).zfill(2)} hs.')
                else:
                    for ii in od_matrix_agg.trip_datetime.dt.date.unique():                    
                        ii = pd.to_datetime(ii)
                        print(f'        {weekDaysMapping[ii.weekday()]} {ii.day} de {monthsMapping[ii.month]} de {ii.year} para las 24 horas del día')
                print('')
                if not only_distance_duration:                    
                    tiempo_estimado = qty_queries * 0.006 
                    
                    if tiempo_estimado > 60:                                         
                        tiempo_estimado = tiempo_estimado / 60
                        tiempo_estimado = f'{int(tiempo_estimado)}:{int(60*(tiempo_estimado - int(tiempo_estimado)) )} horas'
                    else:                        
                        tiempo_estimado = str(int(tiempo_estimado)) + ' minutos'

                    print(f' El tiempo estimado para correr este proceso es de {tiempo_estimado}')
                    print('')


                answer = input("  Ingrese si para continuar ")
                print('')
                if answer.lower() == 'si':
                    print('')
                    print(f'Las consultas van a quedar guardadas en un archivo temporal con el nombre {trips_file}_fecha')
                    print('')
                    
                    coord_origin = od_matrix_agg.head(1)[geo_origin].values[0]

                    for trip_datetime in od_matrix_agg.trip_datetime.unique():

                        trip_datetime = pd.to_datetime(trip_datetime)

                        trips_file_ = current_path / 'tmp' / Path(f'{trips_file}_{str(trip_datetime)[:10]}.csv'.replace(':', '_').replace(' ', '_'))
                        
                        if len(trips_all) > 0:                            
                            trips = trips_all[trips_all.trip_datetime.dt.date == trip_datetime.date()]
                        else:
                            trips = pd.DataFrame([])
                        
                        if not only_distance_duration:

                            trip_datetime_corrected = sup.correct_datetime(trip_datetime, lat=float(coord_origin.split(',')[0]), lon=float(coord_origin.split(',')[1]))

                            n = 0                
                            for _, i in od_matrix_agg[(od_matrix_agg.trip_datetime==trip_datetime)&(od_matrix_agg[geo_origin]!=od_matrix_agg[geo_destination])].iterrows():         

                                trips_new = trip_info_googlemaps(i[geo_origin], 
                                                                 i[geo_destination], 
                                                                 trip_datetime = pd.to_datetime(i.trip_datetime), 
                                                                 trip_datetime_corrected = trip_datetime_corrected,
                                                                 gmaps = gmaps, 
                                                                 transit=transit,
                                                                 driving=driving,
                                                                 walking=walking,
                                                                 bicycling=bicycling)
                                
                                trips_new = trips_new.rename(columns={'geo_origin':geo_origin, 'geo_destination':geo_destination})

                                sk = save_key(key, (transit+((driving)*1.5)+walking+bicycling), current_path)

                                trips = pd.concat([trips,
                                                   trips_new], 
                                                   ignore_index=True)
                                
                                trips_all = pd.concat([trips_all,
                                                   trips_new], 
                                                   ignore_index=True)

                                n += 1

                                if (n % 100) == 0:

                                    n_queries = n*(transit * 1 + driving * 1 + walking * 1 + bicycling*1)
                                    n_queries = '{:,}'.format(n_queries)

                                    print(f'{datetime.datetime.now().strftime("%H:%M:%S")} procesando {n_queries} de {qty_queries}')
                                    trips.to_csv(trips_file_, index=False)

                            trips.to_csv(trips_file_, index=False)

                        else:
                        
                            print(f'Procesando {str(pd.to_datetime(trip_datetime))}')

                            trips_new = trip_info_googlemaps_matrix(od_matrix_agg[(od_matrix_agg.trip_datetime==trip_datetime)&(od_matrix_agg[geo_origin]!=od_matrix_agg[geo_destination])], 
                                                                    geo_origin = geo_origin,
                                                                    geo_destination = geo_destination,
                                                                    trip_datetime = trip_datetime, 
                                                                    gmaps = gmaps,
                                                                    transit = transit, 
                                                                    driving = driving, 
                                                                    walking = walking, 
                                                                    bicycling = bicycling)
                            
                            trips_new = trips_new.rename(columns={'geo_origin':geo_origin, 'geo_destination':geo_destination})
                            
                            sk = save_key(key, (transit+((driving)*1.5)+walking+bicycling)*len(trips_new), current_path)

                            trips = pd.concat([trips,
                                               trips_new], 
                                               ignore_index=True)
                            
                            trips_all = pd.concat([trips_all,
                                                   trips_new], 
                                                   ignore_index=True)

                            trips.to_csv(trips_file_, index=False)

                    print('Proceso finalizado')
                else:
                    trips_all = ''
                    od_matrix_all = pd.DataFrame([])
            else:
                print('')
                print('La fecha/hora del proceso debe ser una fecha/hora mayor a la actual')
                print('')

                trips_all = ''
                od_matrix_all = pd.DataFrame([])


    else:
        
        print(f' Este proceso ya se corrió con anterioridad. Las consultas están guardadas en un archivo temporal con el nombre {trips_file}')
        print(f' Puede borrar este archivo si quiere correr nuevamente el proceso para este mismo día')
        print('')
    
    if len(trips_all) > 0:
        
        od_matrix_all = od_matrix_all.merge(trips_all, 
                                           how='left', 
                                           on=[geo_origin, geo_destination, 'trip_datetime'])
    

    if geo_origin in od_matrix_all.columns: 
        
        for i in od_matrix_all.columns:
            if od_matrix_all[i].dtype == 'float':        
                od_matrix_all.loc[od_matrix_all[geo_origin] == od_matrix_all[geo_destination], i] = od_matrix_all.loc[od_matrix_all[geo_origin] == od_matrix_all[geo_destination], i].fillna(0)

    return od_matrix_all


def trips_gmaps_from_od(origin, 
                        id_origin, 
                        destination, 
                        id_destination, 
                        trip_datetime,                 
                        key, 
                        transit=True,
                        driving=True,
                        walking=False,
                        bicycling=False,
                        normalize=True,
                        res = 8,                        
                        full_day=False,
                        only_distance_duration=False,
                        trips_file = 'trips_file_tmp',
                        current_path=Path()):
    
    '''
    Calcula distancias en modos de transporte entre dos capas cartográficas de origenes y destinos.
    En el caso de que alguna capa sea de polígonos, toma en cuenta el centroide de cada polígono.
    Devuelve una od_matrix de origenes y destinos con las distancias calculadas para cada modo.

    Parámetros
    origin = GeoDataFrame que representa a los puntos de origen
    id_origin = ID del GeoDataFrame origin
    destination = GeoDataFrame que representa a los puntos de destino
    id_destination = ID del GeoDataFrame destino
    trip_datetime = Fecha/hora de la consulta en formato datetime.
    key = key requerida para consultar la API de Google Maps
    transit = Para obtener distancias y tiempos en transporte público (True/False). Por defecto True.
    driving = Para obtener distancias y tiempos en automovil (True/False). Por defecto True.
    walking = Para obtener distancias y tiempos caminando (True/False). Por defecto False.
    bicycling = Para obtener distancias y tiempos en bicicleta (True/False). Por defecto False.
    normalize = Si normalize es True, toma los origenes y destinos desde el centroide del hexágono h3 con la resolución correspondiente. Este permite minimizar la cantidad de consultas de tiempos de viaje en observaciones que están a una distancia muy cercana. Esto permite reducir los costos de consulta en Google Maps.
    res = Resolución de los hexágonos h3 si normalize = True.    
    full_day = Si se realiza consulta para un día completo o para una sola hora del día. Por defecto es False
    only_distance_duration = Si se quiere consultar solo tiempos y distancias en la matriz. La corrida de este proceso es mucho más rápida pero no incluye todas las variables, como distancias de caminata para acceder al transporte público, tiempos en cada modo, etc.
    trips_file = Nombre de archivo temporario para guardar las consultas a Google Maps. Por defecto: 'trips_file_tmp',
    current_path = Directorio de trabajo
    
    Salida: matriz de origenes y destinos con las variables de tiempo y distancias calculadas para los modos solicitados

    '''
    
    if not normalize:
        lat_o, lon_o, lat_d, lon_d = 'lat_o', 'lon_o', 'lat_d', 'lon_d'
        geo_origin, geo_destination = 'origin', 'destination'
    else:
        lat_o, lon_o, lat_d, lon_d = 'lat_o_norm', 'lon_o_norm', 'lat_d_norm', 'lon_d_norm'
        geo_origin, geo_destination = 'origin_norm', 'destination_norm'

    od_matrix_all_ = create_matrix(origin,                                   
                                  destination, 
                                  id_origin = id_origin, 
                                  id_destination = id_destination, 
                                  latlon=False, 
                                  normalize=False, 
                                  duplicates=True,
                                  res = res)
        
    list_trip_datetime = trip_datetime
    if type(trip_datetime) != list:
        list_trip_datetime = [list_trip_datetime]
        
    od_matrix_all = pd.DataFrame([])
    for i in list_trip_datetime:
        od_matrix_all_['trip_datetime'] = i
        od_matrix_all = pd.concat([od_matrix_all, od_matrix_all_], ignore_index=True)
        
    if full_day:
        od_matrix_all_ = od_matrix_all.copy()
        od_matrix_full_day = pd.DataFrame([])
        for n in range(0, 24):                    
            od_matrix_all_['trip_datetime'] = od_matrix_all_.trip_datetime.apply(lambda dt: dt.replace(hour=n))
            od_matrix_all_['hour'] = n
            od_matrix_all = pd.concat([od_matrix_all,
                                       od_matrix_all_], ignore_index=True)
        od_matrix_all = od_matrix_all.sort_values(['trip_datetime', 'hour'])
        
        del od_matrix_all['hour']
        
        od_matrix_all = od_matrix_all.drop_duplicates().reset_index(drop=True)        
        
        
        
    od_matrix = create_matrix(origin,                               
                              destination, 
                              id_origin = id_origin, 
                              id_destination = id_destination, 
                              latlon=False,
                              normalize=normalize,
                              duplicates=False,
                              res = res)

    
    od_matrix = trips_gmaps_process(od_matrix = od_matrix, 
                                    geo_origin=geo_origin,
                                    geo_destination=geo_destination,
                                    trip_datetime = trip_datetime,                 
                                    key = key, 
                                    transit=transit,
                                    driving=driving,
                                    walking=walking,
                                    bicycling=bicycling,
                                    normalize=normalize,
                                    res = res,
                                    trips_file = trips_file,
                                    full_day=full_day,
                                    only_distance_duration=only_distance_duration,
                                    current_path=current_path )
    
    if (not normalize)&(f'{geo_origin}_norm' in od_matrix.columns):
        od_matrix = od_matrix.drop([f'{geo_origin}_norm', f'{geo_destination}_norm'], axis=1)
        od_matrix_all = od_matrix_all.drop([f'{geo_origin}_norm', f'{geo_destination}_norm'], axis=1)
    

    
    if len(od_matrix) > 0:
        
        od_matrix = od_matrix_all.merge(od_matrix, how='left', on=['hex_o', 'hex_d', 'trip_datetime', geo_origin, geo_destination])
        od_matrix['hour'] = od_matrix.trip_datetime.dt.hour
        od_matrix = od_matrix.sort_values(['trip_datetime', 'hour'])
        del od_matrix['hour']
        
    cols = [i for i in [id_origin, id_destination, 'hex_o', 'hex_d', 'origin', 'destination', 'origin_norm', 'destination_norm', 
                        'area_m2', 'cant_pers', 'PCA_1', 'NSE_5', 'NSE_3', 'weight', 'weight%', 
                        'distance_osm_drive', 'distance_osm_walk',
                         'trip_datetime', 'transit_departure_time', 'transit_arrival_time',
                         'transit_distance', 'transit_duration', 'transit_walking_distance', 'transit_walking_duration', 'transit_walking_steps',
                         'transit_transit_distance', 'transit_transit_duration', 'transit_transit_steps', 'transit_walking_distance_origin',
                         'transit_walking_duration_origin', 'driving_distance',
                         'driving_duration', 'driving_duration_in_traffic',
                         'walking_distance', 'walking_duration',
                         'bicycling_distance', 'bicycling_duration'] if i in od_matrix.columns]  
    
    return od_matrix[cols]
    
def trips_gmaps_from_matrix(od_matrix, 
                            trip_datetime,                 
                            key, 
                            geo_origin = 'origin',
                            geo_destination = 'destination',                        
                            transit=True,
                            driving=True,
                            walking=False,
                            bicycling=False,
                            normalize=True,
                            res = 8,
                            trips_file = 'trips_file_tmp',                        
                            current_path=Path()):
    '''
    Realiza consulta a la API de Google Maps a partir de una matriz ya existente. El proceso es similar a trips_gmaps_from_od pero con la diferencia que la matriz ya viene como parámetro.
    
    Parámetros
    od_matrix = GeoDataFrame que representa a la matriz de origenes y destinos
    key = key requerida para consultar la API de Google Maps
    geo_origin = Variable con lat/lon origen. Por defecto: 'origin',
    geo_destination = Variable con lat/lon destino. Por defecto: 'destination',
    trip_datetime = Fecha/hora de la consulta en formato datetime.
    transit = Para obtener distancias y tiempos en transporte público (True/False). Por defecto True.
    driving = Para obtener distancias y tiempos en automovil (True/False). Por defecto True.
    walking = Para obtener distancias y tiempos caminando (True/False). Por defecto False.
    bicycling = Para obtener distancias y tiempos en bicicleta (True/False). Por defecto False.
    normalize = Si normalize es True, toma los origenes y destinos desde el centroide del hexágono h3 con la resolución correspondiente. Este permite minimizar la cantidad de consultas de tiempos de viaje en observaciones que están a una distancia muy cercana. Esto permite reducir los costos de consulta en Google Maps.
    res = Resolución de los hexágonos h3 si normalize = True.        
    trips_file = Nombre de archivo temporario para guardar las consultas a Google Maps. Por defecto: 'trips_file_tmp',
    current_path = Directorio de trabajo
    
    Salida: matriz de origenes y destinos con las variables de tiempo y distancias calculadas para los modos solicitados
    '''

    if normalize:        
        geo_origin = geo_origin + '_norm'        
        geo_destination = geo_destination + '_norm'        

        od_matrix['lat_o_norm'] = od_matrix['hex_o'].apply(suph3.add_geometry, bring='lat').round(5)
        od_matrix['lon_o_norm'] = od_matrix['hex_o'].apply(suph3.add_geometry, bring='lon').round(5)
        od_matrix['lat_d_norm'] = od_matrix['hex_d'].apply(suph3.add_geometry, bring='lat').round(5)
        od_matrix['lon_d_norm'] = od_matrix['hex_d'].apply(suph3.add_geometry, bring='lon').round(5)
        od_matrix[geo_origin] = od_matrix[f'lat_o_norm'].round(5).astype(str) + ', ' + od_matrix[f'lon_o_norm'].round(5).astype(str)
        od_matrix[geo_destination] = od_matrix[f'lat_d_norm'].round(5).astype(str) + ', ' + od_matrix[f'lon_d_norm'].round(5).astype(str)
        od_matrix = od_matrix.drop(['lat_o_norm', 'lon_o_norm', 'lat_d_norm', 'lon_d_norm'], axis=1)


    
    od_matrix_agg = od_matrix.groupby(['hex_o', 'hex_d', geo_origin, geo_destination], as_index=False).size().drop(['size'], axis=1)
    
    od_matrix_agg = trips_gmaps_process(od_matrix = od_matrix_agg, 
                                        geo_origin=geo_origin,
                                        geo_destination=geo_destination,
                                        trip_datetime = trip_datetime,                 
                                        key = key, 
                                        transit=transit,
                                        driving=driving,
                                        walking=walking,
                                        bicycling=bicycling,
                                        normalize=normalize,
                                        res = res,
                                        trips_file = trips_file,
                                        current_path=current_path )
    
    
    
    
    if len(od_matrix_agg) > 0:
        od_matrix = od_matrix.merge(od_matrix_agg, on=['hex_o', 'hex_d', geo_origin, geo_destination])
        

    return od_matrix
    
    
def distances_to_equipments(origin,
                            destination, 
                            id_origin,
                            id_destination,                            
                            trip_datetime,
                            key,                           
                            geo_origin = 'origin',
                            geo_destination = 'destination',
                            processing='pandana',
                            equipement_bring_closest = True,
                            equipment_closest_qty = 2,
                            equipment_type = '',
                            normalize=True,
                            closest_distance=[800, 1500, 2000],
                           current_path=Path()):
    '''
    Calcula las distancias entre una capa geográfica de origenes y un capa de destinos que corresponda a algún tipo de equipamiento (ej. escuelas, hospitales, etc)
    En una primera instancia se calcula una matriz de distancia en OpenStreetMaps desde el centroide de la capa de origen hacia el establecimiento (capa destino)
    Una vez calculadas las distancias para cada origen se obtienen los equipment_closest_qty más cercanos. 
    Una vez seleccionados los más cercanos se realiza la consulta de tiempos de viaje en transporte público o caminando utilizando Google Maps.
    
    Parámetros
    origin = capa de origen (Ejemplo, capa de hexágonos h3)
    destination = capa de destinos
    id_origin = id capa de origen
    id_destination = id capa de destino
    trip_datetime = fecha y hora para realizar la consulta en Google Maps
    key = key para la API de google maps
    geo_origin = Variable con lat/lon origen. Por defecto: 'origin',
    geo_destination = Variable con lat/lon destino. Por defecto: 'destination',
    equipement_bring_closest = Si trae los más cercanos. Por defecto True,
    equipment_closest_qty = Cantidad de equipamientos más cercanos de cada origen. Por defecto 2.
    equipment_type = Si los equipamientos se encuentran categorizados según alguna variable (ej. escuelas primarias y escuelas secundarias). Se van a traer los equipamiento más cercano para cada tipo en forma independiente.
    normalize = Si se quiere utilizar el centroide del hexágono h3 de destino en vez de la ubicación exácta. Esto permite reducir la cantidad y el costo de las consultas a la API de Google Maps.
    closest_distance = Genera una variable con la cantidad de establecimientos en un rango de distancia para cada origen. Por defecto: closest_distance = [800, 1500, 2000] 
    
    Salida
    Matriz de tiempos y distancias de viaje a los establecimientos más cercanos.
    '''

    print('Calcula distancias en Open Street Maps')
    od_matrix_osm = measure_distances_osm(origin = origin, 
                                          id_origin = id_origin, 
                                          destination = destination, 
                                          id_destination = id_destination,           
                                          driving=False,
                                          walking=True,
                                          processing=processing, 
                                          equipement_bring_closest = equipement_bring_closest,
                                          equipment_closest_qty = equipment_closest_qty,
                                          closest_distance = closest_distance,
                                          equipment_type = equipment_type,
                                          current_path=current_path)
    
    print('')
    print('Calcula tiempos en transporte público con Google Maps')
    od_matrix_transit =   trips_gmaps_from_matrix(od_matrix = od_matrix_osm.copy(), 
                                                    trip_datetime = trip_datetime,                 
                                                    key = key, 
                                                    geo_origin = geo_origin,
                                                    geo_destination = geo_destination,                        
                                                    transit=True,
                                                    driving=False,
                                                    walking=False,
                                                    bicycling=False,
                                                    normalize=normalize,
                                                    current_path=current_path)

    print('')
    if 'transit_duration' in od_matrix_transit.columns:
        od_matrix_transit['total_duration'] = od_matrix_transit['transit_duration']

        od_matrix_transit['modo'] = 'walk'
        od_matrix_transit.loc[(od_matrix_transit.transit_transit_duration.notna())&(od_matrix_transit.transit_transit_duration != 0), 'modo'] = 'transit'

        if 'distance_osm_walk' in od_matrix_transit.columns:

            tmp = od_matrix_transit.loc[(od_matrix_transit.transit_duration.notna())&(od_matrix_transit.transit_transit_duration.isna()), :].groupby(['origin_norm', 'destination_norm', 'transit_distance', 'transit_duration'], as_index=False).size().drop(['size'], axis=1)

            tmp['kmh_walk'] = (tmp['transit_distance'] / (tmp['transit_duration'] / 60)).round(2)

            kmh_walk = 4.8
            if len(tmp)>=100:            
                kmh_walk = tmp.kmh_walk.mean().round(1)

            od_matrix_transit['duration_osm_walk'] = (od_matrix_transit['distance_osm_walk'] / kmh_walk * 60).round(1)

            od_matrix_transit.loc[(od_matrix_transit.total_duration.isna())|(od_matrix_transit.total_duration==0), 'total_duration'] = od_matrix_transit['duration_osm_walk']

            if len(equipment_type)==0: equipment_type=[]
            if type(equipment_type)==str: equipment_type=[equipment_type]

            od_matrix_transit = od_matrix_transit.sort_values(equipment_type + [id_origin, 'total_duration'])
            od_matrix_transit['duration_order'] = od_matrix_transit.groupby(equipment_type + [id_origin]).transform('cumcount')
            od_matrix_transit['duration_order'] = od_matrix_transit['duration_order'] + 1

            od_matrix_transit = od_matrix_transit[od_matrix_transit.duration_order==1].reset_index(drop=True)

            var_osm = od_matrix_osm.columns.tolist()
            qty_var = []
            for x in closest_distance:
                var = f'qty_est_{x}m'
                var_osm.remove(var)
                qty_var.append(var)

            if 'distance' in var_osm: var_osm.remove('distance')
            if 'distance_order' in var_osm: var_osm.remove('distance_order')
            if 'duration' in var_osm: var_osm.remove('duration')
            if 'origin' in var_osm: var_osm.remove('origin')
            if 'destination' in var_osm: var_osm.remove('destination')
            if 'distance_osm_walk' in var_osm: var_osm.remove('distance_osm_walk')

            vars_od = [geo_origin, geo_destination]
            if normalize:
                vars_od += [f'{geo_origin}_norm', f'{geo_destination}_norm']

            od_matrix_transit = od_matrix_transit[var_osm+vars_od+['trip_datetime', 'modo', 'distance_osm_walk', 'total_duration']+qty_var]

            od_matrix_transit = od_matrix_transit.rename(columns={'distance_osm_walk': 'distance', 'total_duration':'duration'})

    cols = [i for i in od_matrix_transit.columns if i not in [id_origin, id_destination, 'hex_o', 'hex_d', 'origin', 'destination', 'origin_norm', 'destination_norm']]
    cols = [id_origin, id_destination, 'hex_o', 'hex_d', 'origin', 'destination', 'origin_norm', 'destination_norm'] + cols
    return od_matrix_transit[cols]
    
    
def calculate_green_space(df, city_crs, population, max_distance = [1000, 2000], green_space='', osm_tags={'leisure': ['park', 'playground', 'nature_reserve', 'recreation_ground']} ):
    '''
    Calculo el área de espacio verde y el espacio verde per-cápita en un radio determinado.
    
    df = capa de hexágonos h3. Tiene que contener un campo de población
    city_crs = Proyección adecuada en metros
    population = variable de población de la capa de hexágonos
    max_distance = Distancias a las que se quiere realizar el cálculo. Es el área a partir del centroide del hexágono. Por defecto max_distance = [1000, 2000]
    green_space = Capa geográfica con parques o espacios públicos. De encontrarse esta capa vacia se obtiene de OpenStreetMaps con los parámetros osm_tags={'leisure': ['park', 'playground', 'nature_reserve', 'recreation_ground']} 
    
    Salida
    Capa geográfica con el area y area per capita de espacios verdes en los rangos seleccionados para cada hexágono.
    '''

    if type(max_distance) != list:
        max_distance = [max_distance]

    # Traigo spacios verdes y públicos de OSM
    if len(green_space)==0:       

        green_space = bring_osm(df, tags = osm_tags)
        green_space = green_space[(green_space.geom_type == 'Polygon')|(green_space.geom_type == 'MultiPolygon')]
        green_space[f'green_area_m2'] = green_space.to_crs(city_crs).area.round(1)
        green_space = green_space[green_space[f'green_area_m2']>=5000]

    green_space = gpd.overlay(df[['hex', population, 'geometry']], green_space[['osmid', f'green_area_m2', 'geometry']], how='intersection')
    green_space['osmid_order'] = green_space.groupby(['osmid']).transform('cumcount')
    green_space[f'green_area_m2'] = green_space.to_crs(city_crs).area.round(1)

    for i in max_distance:

        df_buffer = df[['hex', population, 'geometry']].copy()
        df_buffer['geometry'] = df_buffer.to_crs(city_crs).representative_point().buffer(i).to_crs(4326)
        shape = gpd.overlay(df_buffer[['hex', population, 'geometry']], green_space[['osmid', 'osmid_order', f'green_area_m2', 'geometry']], how='intersection')
        shape = shape.rename(columns={f'green_area_m2':f'green_area_m2_in{i}m'})
        
        green_space_pobl = shape.groupby(['osmid', 'osmid_order'], as_index=False).agg({population:'sum', f'green_area_m2_in{i}m':'max'})

        green_space_pobl[f'green_pcapita_m2_in_{i}m'] = round(green_space_pobl[f'green_area_m2_in{i}m'] / green_space_pobl[population], 1)
        green_space_pobl[f'green_area_ha_in{i}m'] = (green_space_pobl[f'green_area_m2_in{i}m'] / 10000).round(1)

        shape = shape[['hex', 'osmid', 'osmid_order', population, 'geometry']].merge(green_space_pobl[['osmid', 'osmid_order', f'green_area_ha_in{i}m', f'green_area_m2_in{i}m', f'green_pcapita_m2_in_{i}m']], on=['osmid', 'osmid_order'], how='left')
        
        shape = shape.groupby('hex', as_index=False).agg({f'green_area_ha_in{i}m':'sum',f'green_area_m2_in{i}m':'sum', f'green_pcapita_m2_in_{i}m':'sum'})

        df_new = df.merge(shape, on='hex', how='left')
        df_new[f'green_pcapita_m2_in_{i}m'] = df_new[f'green_pcapita_m2_in_{i}m'].fillna(0)
        df_new[f'green_area_m2_in{i}m'] = df_new[f'green_area_m2_in{i}m'].fillna(0).astype(int)
        df_new[f'green_area_ha_in{i}m'] = df_new[f'green_area_ha_in{i}m'].fillna(0)

        df = df.merge(df_new[['hex', f'green_area_ha_in{i}m', f'green_area_m2_in{i}m', f'green_pcapita_m2_in_{i}m']])
        
    return df
    

def create_result_dir(current_path=Path()):    
    if not Path(current_path / 'Resultados_files').is_dir(): Path.mkdir(current_path / 'tmp')
    if not Path(current_path / 'Resultados_files').is_dir(): Path.mkdir(current_path / 'Resultados_files')
    if not Path(current_path / 'Resultados_pptx').is_dir(): Path.mkdir(current_path / 'Resultados_pptx')
    if not Path(current_path / 'Resultados_png').is_dir(): Path.mkdir(current_path / 'Resultados_png')
    if not Path(current_path / 'Resultados_pdf').is_dir(): Path.mkdir(current_path / 'Resultados_pdf')
    
    
def indicators_all_day(od_matrix_sample, current_path = Path(), city='', prs=''):
    '''
    Calcula indicadores para una muestra de viajes en automovil que se calcula para el día completo. Puede incluir varios días en el mismo DataFrame (ej. día de semana, sábado y domingo)
    od_matrix_sample = DataFrame con la matriz resultado de los viajes en automovil para las distintas horas de o los días.
    current_path = Directorio de trabajo. Por defecto Path()
    Salida. Dataframe de indicadores. Se guardan los resultados la carpeta de Resultados
    
    '''


    
    create_result_dir(current_path)
    
    if 'driving_duration_in_traffic' in od_matrix_sample.columns:
    
        if len(od_matrix_sample[od_matrix_sample.driving_duration_in_traffic.notna()]) > 0:

            weekDaysMapping = ("Lunes", 
                               "Martes",
                               "Miércoles", 
                               "Jueves",
                               "Viernes", 
                               "Sábado",
                               "Domingo")
            monthsMapping = ("", 
                             "enero", 
                             "febrero", 
                             "marzo", 
                             "abril", 
                             "mayo", 
                             "junio", 
                             "julio", 
                             "agosto", 
                             "septiembre", 
                             "octubre", 
                             "noviembre", 
                             "diciembre")

            od_matrix_sample['kmh'] = (od_matrix_sample['driving_distance'] / (od_matrix_sample['driving_duration_in_traffic'] / 60)).round(2)

            trip_time = od_matrix_sample.groupby('trip_datetime', as_index=False).agg({'driving_duration_in_traffic':'mean', 'kmh': 'mean'})
            trip_time['date'] = trip_time.trip_datetime.dt.date
            trip_time['hour'] = trip_time['trip_datetime'].dt.hour.astype(str).str.zfill(2)

            indicators_all = pd.DataFrame([])
            for i in trip_time.date.unique():

                trip_time_tmp = trip_time[trip_time.date == i]
                _ = ''
                for x in range(0, len(f'{weekDaysMapping[i.weekday()]} {i.day} de {monthsMapping[i.month]}')):
                    _ += '-'

                indicadores = pd.DataFrame([[pd.to_datetime(i).date(), 
                                          f'{weekDaysMapping[i.weekday()]} {i.day} de {monthsMapping[i.month]}',
                                          trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmax()].hour,              
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmax()].driving_duration_in_traffic, 2),
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmax()].kmh, 2),
                                          trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmin()].hour,
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmin()].driving_duration_in_traffic, 2),
                                          round(trip_time_tmp.loc[trip_time_tmp.driving_duration_in_traffic.idxmin()].kmh, 2),    
                                          round(trip_time_tmp.loc[trip_time_tmp.kmh.idxmax()].kmh / trip_time_tmp.loc[trip_time_tmp.kmh.idxmin()].kmh, 2),
                                          round(trip_time_tmp.driving_duration_in_traffic.mean(), 2),
                                          round(trip_time_tmp.kmh.mean(), 2)
                                          ]], 

                                         columns=['Date', 
                                                  'Detalle día',
                                                 'Hora Punta',                      
                                                 'Tiempo de viaje en hora punta (min)',
                                                 'Velocidad de viaje en hora punta (kmh)',
                                                 'Hora Valle', 
                                                 'Tiempo de viaje en hora valle (min)',
                                                 'Velocidad de viaje en hora valle (kmh)',
                                                 'Índice de congestión',
                                                 'Tiempo promedio de los viajes (min)', 
                                                 'Velocidad promedio de los viajes (kmh)', ])

                indicators_all = pd.concat([indicadores, indicators_all], ignore_index=True)

            df = indicators_all.set_index('Detalle día').T.reset_index().rename(columns={'index':'Detalle día'}).copy()

            fig = Figure(figsize=(13.5,13.5), dpi=160)
            canvas = FigureCanvas(fig)
            ax = fig.add_subplot(111)

            #hide the axes
            fig.patch.set_visible(False)
            ax.axis('off')
            # ax.axis('tight')
            table = ax.table(cellText=df.values, colLabels=df.columns, loc='center', colLoc='center', cellLoc='center')
            # table.set_fontsize(18)
            # table.scale(1,3)
            #display table
            fig.tight_layout()
            fig.savefig(current_path / 'Resultados_png' / f'{city}_indicadores_dia_completo.png', dpi=300)
            fig.savefig(current_path / 'Resultados_pdf' / f'{city}_indicadores_dia_completo.pdf', dpi=300)


            indicadores.to_csv(current_path / 'Resultados_files' / 'indicadores_dia_completo.csv', index=False)

            with sns.axes_style('darkgrid', {"axes.facecolor": "#d4dadc", 'figure.facecolor': "#d4dadc"}):

                # Tiempos promedios
                fig = Figure(figsize=(6,3), dpi=100)
                canvas = FigureCanvas(fig)
                ax = fig.add_subplot(111)


                for i in trip_time.date.unique():
                    var = f'{weekDaysMapping[i.weekday()]} {i.day}/{i.month}/{i.year}'
                    current_day = trip_time.loc[trip_time.date==i, ['hour', 'driving_duration_in_traffic']].reset_index(drop=True).rename(columns={'driving_duration_in_traffic':var})        
                    current_day.plot(ax=ax, legend=True, lw=.6)


                ax.set_title(f'Tiempos promedio de viaje', fontsize=8)
                ax.set_xlabel('Hora', fontsize=8)
                ax.set_ylabel('Tiempo promedio\n(minutos)', fontsize=8)
                ax.set_xticks(list(range(0, 24)))
                ax.tick_params(labelsize=6);

                box = ax.get_position()
                ax.set_position([box.x0, box.y0, box.width * 0.87, box.height])
                ax.legend(loc='center left', bbox_to_anchor=(1, 0.5), fontsize=6)
                fig.savefig(current_path / 'Resultados_png'  / f'{city}_tiempos_dia_completo.png', dpi=300)
                fig.savefig(current_path / 'Resultados_pdf'  / f'{city}_tiempos_dia_completo.pdf', dpi=300)
                if prs == '': display(fig)


                # Velocidades promedio
                fig = Figure(figsize=(6,3), dpi=100)
                canvas = FigureCanvas(fig)
                ax = fig.add_subplot(111)


                for i in trip_time.date.unique():
                    var = f'{weekDaysMapping[i.weekday()]} {i.day}/{i.month}/{i.year}'
                    current_day = trip_time.loc[trip_time.date==i, ['hour', 'kmh']].reset_index(drop=True).rename(columns={'kmh':var})        
                    current_day.plot(ax=ax, legend=True, lw=.6)

                ax.set_title(f'Velocidad promedio de viaje', fontsize=8)
                ax.set_xlabel('Hora', fontsize=8)
                ax.set_ylabel('Velocidad promedio\n(kmh)', fontsize=8)
                ax.set_xticks(list(range(0, 24)))

                ax.tick_params(labelsize=6);

                box = ax.get_position()
                ax.set_position([box.x0, box.y0, box.width * .87, box.height])
                ax.legend(loc='center left', bbox_to_anchor=(1, 0.5), fontsize=6)
                fig.savefig(current_path / 'Resultados_png' / f'{city}_kmh_dia_completo.png', dpi=300)
                fig.savefig(current_path / 'Resultados_pdf' / f'{city}_kmh_dia_completo.pdf', dpi=300)
                if prs == '': display(fig)

                if prs != '':
                    slide = pptx_addtitle(prs=prs, slide='',  title='', left=14.5, top=0, width=9, new=True, fontsize=48)    
                    pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_tiempos_dia_completo.png',  left=.5, top=1.5, width=11)
                    pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_kmh_dia_completo.png'    ,  left=11.5, top=1.5, width=11)

                    pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_indicadores_dia_completo.png',  left=3, top=8, width=18)
                else:
                    display(indicators_all.set_index('Detalle día').T)

        else:
            print('No hay datos para procesar indicadores')
            indicators_all = pd.DataFrame([])
    else:
            print('No hay datos para procesar indicadores')
            indicators_all = pd.DataFrame([])
    
    return indicators_all
    
    
    

    
    
def calculate_nse_in_hexagons(censo,
                              id_censo = 'RADIO_LINK',                          
                              population='cant_pers',
                              vars_nse = '', 
                              city_crs = '',
                              current_path = Path(),
                              city='',
                              res=8,
                              run_always=True):
    
    create_result_dir(current_path=current_path)
    
    if (not Path(current_path / 'Resultados_files' / f'{city}_hexs.geojson').is_file())|(run_always):


        censo = calculate_nse(censo, 
                              vars_nse, 
                              population='cant_pers', 
                              show_map=False)

        hexs = suph3.create_h3(censo, 
                         res=res, 
                         show_map=False)

        hexs = distribute_population(gdf=censo, 
                                     id_gdf='RADIO_LINK', 
                                     hexs=hexs, 
                                     id_hexs='hex', 
                                     population='cant_pers', 
                                     pca='PCA_1', 
                                     crs=city_crs, 
                                     q=[5, 3],
                                     order_nse = [['Alto', 'Medio-Alto', 'Medio', 'Medio-Bajo', 'Bajo'],
                                                  ['Alto', 'Medio', 'Bajo']],
                                     show_map=True)

        hexs.to_file(current_path / 'Resultados_files' / f'{city}_hexs.geojson')
        print('')
        print('Se guardó el archivo hexs.geojson en', current_path / f'{city}_hexs.geojson')
        print('')
    else:
        hexs = gpd.read_file(current_path / 'Resultados_files' / f'{city}_hexs.geojson')

    return hexs
    
    
def calculate_activity_density(hexs,
                               tags = {'amenity':True},
                               cantidad_clusters = 8,
                               city_crs = '',
                               current_path = Path(),
                               city='',                              
                               run_always=True):
    
    create_result_dir(current_path=current_path)
    
    if (not Path(current_path / 'Resultados_files' / f'{city}_activity_density.geojson').is_file())|(run_always):


        amenities = bring_osm(hexs, tags = tags)
        amenities = assign_weights(amenities)
        densidad_actividad_result, scores, amenities2 = activity_density(amenities, 
                                                                  city_crs, 
                                                                  cantidad_clusters = cantidad_clusters,                                                               
                                                                  show_map = True)

        densidad_actividad_result.to_file(current_path / 'Resultados_files' / f'{city}_activity_density.geojson')
        print('')
        print('Se guardó el archivo hexs.geojson en', current_path / f'{city}_activity_density.geojson')
        print('')
    else:
        densidad_actividad_result = gpd.read_file(current_path / 'Resultados_files' / f'{city}_activity_density.geojson')

    return densidad_actividad_result
    
    
def calculate_od_matrix_all_day(origin, 
                                id_origin, 
                                destination, 
                                id_destination,                                 
                                trip_datetime, 
                                key, 
                                samples_origin = 12,
                                samples_destination = 6,
                                transit=False,
                                driving=True,
                                walking=False,
                                bicycling=False,
                                current_path=Path(), 
                                city = '',
                                normalize=False,
                                run_always=True):
    
    create_result_dir(current_path=current_path)
    


    if (not (current_path / 'tmp' / f'{city}_origin_sample.geojson').is_file())|(run_always):    
        origin_sample = hexs.sample(samples_origin, random_state=2).copy()
        origin_sample.to_file(current_path / 'tmp' / f'{city}_origin_sample.geojson')
        destination_sample = densidad_actividad.sample(samples_destination, random_state=2).copy()
        destination_sample.to_file(current_path / 'tmp' / f'{city}_destination_sample.geojson')
    else:    
        origin_sample =  gpd.read_file(current_path / 'tmp' / f'{city}_origin_sample.geojson')
        destination_sample =  gpd.read_file(current_path / 'tmp' / f'{city}_destination_sample.geojson')


    od_matrix_all_day = trips_gmaps_from_od(origin = origin_sample, 
                                            id_origin = id_origin, 
                                            destination = destination_sample, 
                                            id_destination = id_destination, 
                                            trip_datetime = trip_datetime,                 
                                            key = key, 
                                            transit=transit,
                                            driving=driving,
                                            walking=walking,
                                            bicycling=bicycling,
                                            full_day=True,
                                            normalize=normalize,
                                            only_distance_duration=True,
                                            current_path=current_path)


    if len(od_matrix_all_day) > 0:
        od_matrix_all_day.to_csv(current_path / 'Resultados_files' / f'{city}_od_matrix_all_day.csv', index=False)
        print('')
        print('Se guardó el archivo od_matrix_all_day.geojson en', current_path / f'{city}_od_matrix_all_day.csv')
        print('')
    
    if len(od_matrix_all_day) > 0:
        indicators_all = indicators_all_day(od_matrix_all_day, current_path, city=city)
    
    return od_matrix_all_day
    
    
def crop_imagen(filePath, reduce=1, altura_max=0, ancho_max=0, save=True, crop_left = 0, crop_top = 0, crop_right = 0, crop_bottom = 0):
    
    # Trim all png images with white background in a folder
    # Usage "python PNGWhiteTrim.py ../someFolder padding"

    image=Image.open(filePath)
    image.load()
    imageSize = image.size #tuple
    
    ## QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    # remove alpha channel
    invert_im = image.convert("RGB")
    # invert image (so that white is 0)
    invert_im = ImageOps.invert(invert_im)
    imageBox = invert_im.getbbox()
    cropped=image.crop(imageBox)
    ## FIN DE QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    
    #REDUCE TAMAÑO
    _size=[]
    # calculates percentage to reduce image by maintaining proportion
    if altura_max>0: _size.append((altura_max/(cropped.height/38)))
    if ancho_max>0: _size.append((ancho_max/(cropped.width/38)))
    if len(_size) > 0: reduce = min(_size)
    
    if reduce < 1:
        basewidth = int(cropped.width * reduce)
        wpercent = (basewidth/float(cropped.size[0]))
        hsize = int((float(cropped.size[1])*float(wpercent)))
        # cropped.resize actually does the resizing
        cropped = cropped.resize((basewidth,hsize), Image.ANTIALIAS)
    
    if crop_left + crop_top + crop_right + crop_bottom > 0:
        width, height = cropped.size 
        crop_right = width - crop_right 
        crop_bottom = height - crop_bottom            
        cropped=cropped.crop((crop_left, crop_top, crop_right, crop_bottom))

    # save the image as cropped
    if save:
        filePath = filePath[0: filePath.find('.')]+'_cropped'+filePath[filePath.find('.'):len(filePath)]
        cropped.save(filePath)
        return filePath
    else:
        return cropped

def pptx_addtitle(prs, slide='', title='', top=0, left=0, width=10, height=1, new=True, fontsize=24, fontcolor='blue', bold=True):

    blank_slide_layout = prs.slide_layouts[6] # Using layout 6 (blank layout)
    # if new create blank slide
    if new:
        slide = prs.slides.add_slide(blank_slide_layout)

    # Set the slides background colour
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(212, 218, 220) # RGBColor(212, 218, 220) is the color of water on the contextily tiles

    # translates from cm to inches
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    # adds a text box onto the slide object
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Gill Sans'
    p.font.color.rgb = RGBColor(64,64,64) # (105,105,105) CSS Dim Grey
    if bold is True:
        p.font.bold = True
    p.font.size = Pt(fontsize)
    p.alignment = PP_ALIGN.LEFT
    #p.font.color = fontcolor
    # many more parameters available

    return slide

def pptx_addtext(prs, slide='', text='', top= 0, left=0, width=10, height=1):
    blank_slide_layout = prs.slide_layouts[6]
    
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.title = text
    # p.font.bold = True
    
    p.alignment = PP_ALIGN.RIGHT
    
    return slide

def pptx_addpic(prs, slide, img_path,  left=0, top=0, width=0, altura_max=0, ancho_max=0, crop_left = 0, crop_top = 0, crop_right = 0, crop_bottom = 0):
    # for adding all maps and graphs
    # altura_max and ancho_max in cm
    blank_slide_layout = prs.slide_layouts[6]

    img_path = str(img_path)

    if os.path.exists(img_path):
        # crop_imagen crops the image
        # NB commented out 20200514
        img_path = crop_imagen(img_path, reduce=1, altura_max=altura_max, ancho_max=ancho_max, save=True, crop_left=crop_left, crop_top=crop_top, crop_right=crop_right, crop_bottom=crop_bottom)
        
        # control position
        left = Inches(left)
        top = Inches(top)
        width  = Inches(width)
        # add to the slide
        if width!=0:
            slide_return = slide.shapes.add_picture(img_path, left, top, width) 
        else:
            slide_return = slide.shapes.add_picture(img_path, left, top) 
        
        os.remove(img_path)
        
        return slide_return
        
        
def print_density_nse(hexs, 
                      population = 'cant_pers', 
                      k = 4, 
                      nse = 'NSE_5', 
                      current_path = Path(), 
                      city = '',
                      prs = ''):
    
    create_result_dir(current_path)
    
    hexs['density_ha'] = round(hexs[population] / (hexs['area_m2']/10000),1).round().astype(int)

    bins = [0] + mapclassify.NaturalBreaks(hexs['density_ha'], k=k).bins.tolist()
    bins_labels = [f'{int(bins[n])} - {int(bins[n+1])}' for n in range(0, len(bins)-1)]
    hexs['density_ha_label'] = pd.cut(hexs.density_ha, bins=bins, labels=bins_labels)

    # Find the centre of the reprojected zonas
    w,s,e,n = hexs.to_crs(3857).total_bounds
    cx = (e+w)/2 # centre point x
    cy = (n+s)/2 # centre point y

    # For the actual plot define a crop that is tighter to minimise unused space at the edge of the map
    crop_extent =  max(abs((e-w)/2), abs((n-s)/2)) * .90 # plus 5%
    crop_w = cx-crop_extent
    crop_s = cy-crop_extent
    crop_e = cx+crop_extent
    crop_n = cy+crop_extent

    fig = Figure(figsize=(13.5,13.5), dpi=40)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)

    hexs.to_crs(3857).plot(ax=ax, column='density_ha_label', categorical = True, lw=0, alpha=.6, cmap='PuBuGn', legend=True,
                          legend_kwds={'loc': 'best', 'frameon': True, 'edgecolor':'white', 'facecolor':None, "title":'Densidad\n(Pers/ha)', 'title_fontsize':14, 'fontsize':14})

    ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=6)

    # Manual plot settings
    ax.set_xlim(crop_w,crop_e)
    ax.set_ylim(crop_s,crop_n)

    ax.set_title('Densidad Poblacional\nPers/ha', fontsize=18)
    ax.axis('off');

    fig.tight_layout()
    fig.savefig(current_path / 'Resultados_png'  / f'{city}_density_ha.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf'  / f'{city}_density_ha.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig) 
        
    hexs['NSE_X'] = hexs[nse].replace({'1 - Alto':'Alto', '2 - Medio-Alto':'Medio-Alto', '3 - Medio':'Medio',  '4 - Medio-Bajo':'Medio-Bajo', '5 - Bajo':'Bajo'})
    hexs['NSE_X'] = pd.CategoricalIndex(hexs['NSE_X'], categories= ['Alto', 'Medio-Alto', 'Medio', 'Medio-Bajo', 'Bajo'])
        
    
    fig = Figure(figsize=(13.5,13.5), dpi=40)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)

    hexs.to_crs(3857).plot(ax=ax, column='NSE_X', categorical = True, lw=0, alpha=.6, cmap='YlOrRd', legend=True,
                          legend_kwds={'loc': 'best', 'frameon': True, 'edgecolor':'white', 'facecolor':None, "title":'Nivel Socioeconómico', 'title_fontsize':14, 'fontsize':14})
    ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=6)
    
    del hexs['NSE_X']
    
    # Manual plot settings
    ax.set_xlim(crop_w,crop_e)
    ax.set_ylim(crop_s,crop_n)
    
    ax.set_title('Nivel Socioeconómico', fontsize=18)
    ax.axis('off');
    fig.tight_layout()
    fig.savefig(current_path / 'Resultados_png'  / f'{city}_NSE.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf'  / f'{city}_NSE.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig) 
        
    if prs != '':
        slide = pptx_addtitle(prs=prs, slide='',  title='', left=14.5, top=0, width=9, new=True, fontsize=48)    
        pptx_addpic(prs=prs, slide=slide, img_path= current_path / 'Resultados_png'  / f'{city}_NSE.png',          left=.5, top=1.5, width=11)
        pptx_addpic(prs=prs, slide=slide, img_path= current_path / 'Resultados_png'  / f'{city}_density_ha.png',  left=12.5, top=1.5, width=11)
    
    
        
def calculate_avg_time_distance(hexs, 
                                od_matrix,                                 
                                population='cant_pers'):
    
    indicators_vars = [i for i in od_matrix.columns if ('time' not in i)&(('transit_' in i)|('driving_' in i)|('osm_' in i)|('walking_' in i)|('bicycling_' in i))]

    common_vars = ['area_m2', population, 'PCA_1', 'NSE_5', 'NSE_3', 'trip_datetime']

    od_matrix_agg = od_matrix.groupby('hex')[common_vars].first().reset_index()

    for var in indicators_vars:
        val = od_matrix[od_matrix[var].notna()].hex.value_counts().value_counts(normalize=True).reset_index().rename(columns={'index':'qty'})
        val['hex_agg'] = val.hex.cumsum()
        for i, row in val.iterrows():
            if row.hex_agg > .85:
                qty_val = row.qty
                break

        val = od_matrix[od_matrix[var].notna()].hex.value_counts().reset_index().rename(columns={'index':'hex', 'hex':'qty_val'})
        val = val[val.qty_val >= qty_val]
        od_matrix_val = od_matrix.loc[od_matrix.hex.isin(val.hex.tolist()), ['hex', var, 'weight']].reset_index(drop=True)
        
        od_matrix_val = od_matrix_val.groupby('hex').apply(lambda x: np.average(x[var], weights=x['weight'])).round(2).reset_index().rename(columns={0:var})
        od_matrix_agg = od_matrix_agg.merge(od_matrix_val, how='left', on='hex')


    downtown = od_matrix.loc[od_matrix.weight == od_matrix.weight.max(), ['hex']+indicators_vars]
    for i in downtown.columns:
        if i != 'hex':
            downtown = downtown.rename(columns = {i: i+'_downtown'})

    indicators_vars = indicators_vars + downtown.columns[1:].tolist()

    od_matrix_agg = od_matrix_agg.merge(downtown, how='left', on='hex')

    od_matrix_agg = hexs[['hex', 'geometry']].merge(od_matrix_agg, on='hex')
    
    return od_matrix_agg


def print_graphs(   od_matrix_new,
                    od_matrix_new_w,
                    var,                     
                    colors_dict = '',                    
                    k = 5,
                    population = 'cant_pers',
                    equipment_type = '',
                    equipment_type_title = '',
                    alpha=.6,
                    current_path = Path(),
                    city = '',
                    prs = ''):
    
# if True:
#     od_matrix_new = df.copy()
#     od_matrix_new_w = df_w.copy()
#     equipment_type_title = ''
    


    colors = {  'distance_osm_drive': 'YlGn',
                'distance': 'YlGn',
                'transit_distance': 'YlGn',
                'transit_walking_distance': 'YlGn',
                'transit_transit_distance': 'YlGn',
                'transit_walking_distance_origin': 'YlGn',
                'transit_distance_downtown': 'YlGnBu',
                'transit_walking_distance_downtown': 'YlGnBu',
                'transit_transit_distance_downtown': 'YlGnBu',
                'transit_walking_distance_origin_downtown': 'YlGnBu',
                'driving_distance': 'YlGn',
                'bicycling_distance': 'YlGn',
                'distance_osm_walk': 'YlGn',
                'distance_osm_walk_downtown': 'YlGnBu',
                'walking_distance': 'YlGn',
                'distance_osm_drive_downtown': 'YlGnBu',
                'driving_distance_downtown': 'YlGnBu',
                'transit_duration': 'PuBuGn',
                'duration': 'PuBuGn',
                'transit_walking_duration': 'PuBuGn',
                'transit_walking_duration_origin': 'PuBuGn',
                'transit_duration_downtown': 'YlOrBr',
                'transit_walking_duration_downtown': 'YlOrBr',
                'transit_walking_duration_origin_downtown': 'YlOrBr',
                'driving_duration_in_traffic': 'PuBuGn',
                'driving_duration_in_traffic_downtown': 'YlOrBr',
                'walking_duration': 'PuBuGn',
                'bicycling_duration': 'PuBuGn',
                'transit_transit_steps': 'YlOrRd',
                'transit_transit_steps_downtown': 'YlOrRd' ,
                'transit_transit_steps_downtown': 'Greens' }
    
    try:
        cmap = colors_dict[var]
    except:
        try:
            cmap = colors[var]
        except:
            cmap = 'YlGn'
    
    desc = ''
    label = ''
    if 'distance' in var: desc += 'Distancias'
    if 'duration' in var: desc += 'Tiempos'
    if 'steps' in var: desc += 'Etapas'
    if 'transit' in var: desc += ' en transporte público'
    if 'driving' in var: desc += ' en automóvil'
    if 'walking' in var: desc += ' caminando'
    if 'bicycling' in var: desc += ' en bicicleta'
    if 'transit_walking_distance' in var: desc = 'Distancias a la parada'
    if 'transit_walking_distance_origin' in var: desc = 'Distancias a la parada en origen'
    if 'transit_walking_duration' in var: desc = 'Tiempos caminando a la parada'
    if 'green_area_ha_in2000m' in var: desc = 'Hectáreas de áreas verdes per capita en 2000m'

    if 'downtown' in var: desc += '\n(viajes al centro)'

    if 'distance' in var: label = 'kms'
    if 'duration' in var: label = 'minutos'
    if 'steps' in var: label = 'steps'
    if 'walking_distance' in var: label = 'mts'
    if 'green_area_ha_in' in var: label = 'mts'

    if 'distance' in var: label_short = 'Distancia (kms)'
    if 'duration' in var: label_short = 'Tiempos de viaje (minutos)'
    if 'steps' in var: label_short = 'Modos utilizados'
    if 'walking_distance' in var: label_short = 'Distancia (mts)'
    if 'green_area_ha_in' in var: label_short = 'Áreas Verdes (ha)'

    if 'Distancias' in desc:
        legtitle = 'Distancias (kms)'
    if 'Tiempos' in desc:
        legtitle = 'Tiempos (min)'
    if 'Etapas' in desc:
        legtitle = 'Etapas'
    if 'áreas verdes' in desc:
        legtitle = 'Áreas Verdes (ha)'

    od_matrix_new = od_matrix_new[od_matrix_new[var].notna()].copy()
    od_matrix_new_w = od_matrix_new_w[od_matrix_new_w[var].notna()].copy()
        
    if equipment_type_title != 'x':        
        if var == 'distance':
            legtitle = 'Distancias (mts)'
            od_matrix_new['distance'] = od_matrix_new['distance'] * 1000
            od_matrix_new_w['distance'] = od_matrix_new_w['distance'] * 1000
            desc = desc+' '+equipment_type_title    
        if var == 'duration':
            desc = 'Tiempos de viaje - ' + equipment_type_title + '\nTransporte público o caminando'
                
            
        

    if 'walking_distance' in var:
        od_matrix_new[var] = od_matrix_new[var] * 1000
        od_matrix_new_w[var] = od_matrix_new_w[var] * 1000

    create_result_dir(current_path)    
    
    ### Imprimo Mapa

    # Find the centre of the reprojected zonas
    w,s,e,n = od_matrix_new.to_crs(3857).total_bounds
    cx = (e+w)/2 # centre point x
    cy = (n+s)/2 # centre point y

    # For the actual plot define a crop that is tighter to minimise unused space at the edge of the map
    crop_extent =  max(abs((e-w)/2), abs((n-s)/2)) * 1.05 # plus 5%
    crop_w = cx-crop_extent
    crop_s = cy-crop_extent
    crop_e = cx+crop_extent
    crop_n = cy+crop_extent
    
    fig = Figure(figsize=(13.5,13.5), dpi=35)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)
    
    plt.rcParams.update({"axes.facecolor": '#d4dadc', 'figure.facecolor': '#d4dadc'})   

    od_matrix_new.to_crs(3857).plot(ax=ax, 
                                    column=var, 
                                    alpha=alpha, 
                                    lw=0, 
                                    categorical = True, 
                                    legend=False, 
                                    cmap = cmap, 
                                    scheme='Quantiles', 
                                    k=k)

    ax.set_title(desc, fontsize=18, weight='bold')
    ctx.add_basemap(ax, source=ctx.providers.CartoDB.Positron, attribution_size=3)

    # Manual plot settings
    ax.set_xlim(crop_w,crop_e)
    ax.set_ylim(crop_s,crop_n)

    ctx.add_attribution(ax, "@ CARTO @OpenStreetMap contributors", font_size=10)
    ax.axis('off')

    if not 'steps' in var:
        bins = mapclassify.Quantiles(od_matrix_new[var], k=k).bins.tolist()
    else:
        bins = mapclassify.Quantiles(od_matrix_new[var], k=3).bins.tolist()

    # colorbar
    cax_position = "top"
    cax_size = "1.5%"
    cax_pad = "-2%"
    cax_aspect = 0.045 #ancho

    cmap = mpl.cm.get_cmap(cmap)
    norm = mpl.colors.BoundaryNorm(bins, cmap.N, extend='min')

    # create an axes on the side of ax. 
    divider = make_axes_locatable(ax)
    cax = divider.append_axes(cax_position, size=cax_size, pad=cax_pad, aspect=cax_aspect)
    mpl.colorbar.ColorbarBase(cax, cmap=cmap, norm=norm, spacing='uniform', orientation='horizontal', label=label,format=tick.FormatStrFormatter('%.0f'), alpha=alpha)

    fig.tight_layout()

    fig.savefig(current_path / 'Resultados_png' / f'{city}_{var}_map.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf' / f'{city}_{var}_map.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig)

    for x in range(0, len(od_matrix_new_w['NSE_5'].unique())):
        od_matrix_new_w['NSE_5'] = od_matrix_new_w['NSE_5'].str.replace(f'{x+1} - ', '')

    ### Imprimo boxplot
    
    # Create figure, canvas, axis
    fig = Figure(figsize=(7,4.5), dpi=70)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)

    plt.rcParams.update({"axes.facecolor": '#d4dadc', 'figure.facecolor': '#d4dadc'})   

    sns.boxplot(x='NSE_5', y=var, data=od_matrix_new_w, palette='PuOr', order=['Bajo', 'Medio-Bajo', 'Medio', 'Medio-Alto', 'Alto'], ax=ax) 

    ax.set_title(label=desc, fontsize=16, weight='bold')
    ax.tick_params(labelsize=14)
    ax.set_xlabel('Nivel socioeconómico', fontsize=14)
    ax.set_ylabel(ylabel=label_short, fontsize=14)


    fig.savefig(current_path / 'Resultados_png' / f'{city}_{var}_boxplot.png', facecolor='#d4dadc', dpi=300);
    fig.savefig(current_path / 'Resultados_pdf' / f'{city}_{var}_boxplot.pdf', facecolor='#d4dadc', dpi=300);
    if prs == '': display(fig)
            
    
    ### Imprimo gráfico de porcentajes de la población
    
    bins_labels = []
    lower = 0
    for i in bins:        
        
        if 'steps' not in var: 
            upper = int(round(i))            
        else:            
            upper = round(i, 2)
            
        bins_labels = bins_labels + [f'{lower} - {upper}']
        lower = upper

    od_matrix_new_w['bins_w'] = pd.cut(od_matrix_new_w[var], [0]+bins, labels = bins_labels)
    od_matrix_new_w['NSE_w'] = od_matrix_new_w['NSE_5'].replace({'Alto': '5 - Alto', 'Medio-Alto':'4 - Medio-Alto', 'Medio':'3 - Medio', 'Medio-Bajo':'2 - Medio-Bajo', 'Bajo':'1 - Bajo'})
    
    # Create figure, canvas and axis
    fig = Figure(figsize=(7,4.5), dpi=70)
    canvas = FigureCanvas(fig)
    ax = fig.add_subplot(111)
    
    plt.rcParams.update({"axes.facecolor": '#d4dadc', 'figure.facecolor': '#d4dadc'})   

    dfw = (pd.crosstab(index=od_matrix_new_w['NSE_w'],
                        columns=od_matrix_new_w['bins_w'], 
                        values=od_matrix_new_w[population], 
                        aggfunc='sum', 
                        normalize='index').fillna(0)) * 100
    dfw.plot(kind='bar', 
              stacked=True, 
              cmap=cmap,
              ax=ax)

    # Manually adjust ax settings
    ax.set_title(desc, fontsize=15, weight='bold')

    ax.tick_params(labelsize=15)
    ax.xaxis.set_tick_params(rotation=0)
    ax.set_xlabel('Nivel Socioeconómico', fontsize=12, rotation='horizontal')
    ax.set_ylabel('Porcentaje de la población', fontsize=12)

    box = ax.get_position()
    ax.set_position([box.x0, box.y0, box.width * 0.75, box.height])
    # Put a legend to the right of the current axis
    ax.tick_params(axis='both', which='major', labelsize=8)
    ax.legend(loc='center left', bbox_to_anchor=(1, 0.5), fontsize=8, title=legtitle, title_fontsize=8, frameon='white')

    # Save the figure, specifying the figure background colour
    fig.subplots_adjust(right=0.75, bottom=0.4)
    fig.tight_layout()

    labels = [i.get_text() for i in ax.get_xticklabels()]
    labels = [i[4:] for i in labels]
    ax.set_xticklabels(labels)

    # labels = [i.get_text() for i in ax.get_yticklabels()]
    # labels_new = [f'{i}%'.format() for i in labels]

#     label_format = '{:.0%}'
# #     ticks_loc = ax.get_yticks().tolist()
#     ticks_loc = [0, .2, .4, .6, .8, 1, 1]
#     ax.set_yticks(ax.get_yticks().tolist())
#     ax.set_yticklabels([label_format.format(x) for x in ticks_loc])

    fig.savefig(current_path / 'Resultados_png' / f'{city}_{var}_bar.png', facecolor='#d4dadc', dpi=300)
    fig.savefig(current_path / 'Resultados_pdf' / f'{city}_{var}_bar.pdf', facecolor='#d4dadc', dpi=300)
    if prs == '': display(fig)

    if prs != '':
        slide = pptx_addtitle(prs=prs, slide='',  title='', left=13.5, top=0, width=10, new=True, fontsize=48)

        pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_{var}_map.png',  left=0, top=0, width=13)
        pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_{var}_boxplot.png', left=14, top=1, width=9)
        pptx_addpic(prs=prs, slide=slide, img_path=current_path / 'Resultados_png' / f'{city}_{var}_bar.png', left=14.5, top=7, width=9)


def print_time_distance(hexs, 
                        od_matrix, 
                        indicators_vars='', 
                        equipment_type='',
                        colors_dict = '',
                        current_path=Path(), 
                        population = 'cant_pers',
                        k = 5,
                        alpha=.6,
                        city='', 
                        prs=''):
    
    if len(indicators_vars) == 0: indicators_vars=od_matrix.columns
    
    indicators_vars = [i for i in od_matrix.columns if (i in indicators_vars)&
                                                      (('transit_' in i)|
                                                       ('driving_' in i)|
                                                       ('osm_' in i)|
                                                       ('walking_' in i)|
                                                       ('bicycling_' in i)|
                                                       ('distance' == i)|
                                                       ('distance' == i)|
                                                       ('duration' == i)|
                                                       ('green_area_ha_in' in i))]
    
    
    if len(equipment_type) == 0:
        od_matrix['equipment_type_tmp'] = 'x'
        equipment_type = ['equipment_type_tmp']

    if type(equipment_type) == str:
        equipment_type = [equipment_type]

    for _, i in od_matrix.groupby(equipment_type).size().reset_index().iterrows():    
        df = od_matrix[(od_matrix[equipment_type] == i[equipment_type]).all(axis=1)].copy()
        
        df = hexs[['hex', 'geometry']].merge(df)
    
        df_w = sup.reindex_df(df, weight_col = population, div=10)
        
        equipment_type_title = i[equipment_type].values[0]
        for x in i[equipment_type].values[1:].tolist():
            equipment_type_title += ' - '+x
        
        for var in indicators_vars:
            
            print_graphs(df,
                         df_w, 
                         var,
                         colors_dict = colors_dict,
                         population = population,
                         equipment_type = equipment_type,
                         equipment_type_title = equipment_type_title,
                         k = k,
                         alpha=alpha,
                         current_path = current_path,
                         city = city,
                         prs = prs)
                         
                         
def create_pptx(hexs, 
                od_matrix_all_day='',
                od_matrix='',
                od_establecimientos='',
                hexs_green_space = '',
                equipment_type='',
                current_path = Path(), 
                city=''):

    prs = Presentation()
    prs.slide_height = Inches(13.5)
    prs.slide_width = Inches(24)

    file_pptx = current_path / 'Resultados_pptx' / f'{city}_Accesibilidad.pptx'

    print('Densidad y nivel socioeconómico')
    print_density_nse(hexs, 
                      current_path = current_path, 
                      city=city,
                      prs=prs)

    if len(od_matrix_all_day) > 0:
        print('Índicadores de día completo')
        indicators_all = indicators_all_day(od_matrix_all_day, current_path, city=city, prs=prs)

    if len(od_matrix) > 0:
        print('Isocronas de tiempos y distancias')
        indicators_vars = ['distance_osm_walk',
                           'transit_duration',
                           'driving_duration_in_traffic',
                           'transit_walking_distance', 
                           'transit_walking_distance_origin',                    
                           'walking_distance', 
                           'walking_duration', 
                           'bicycling_distance', 
                           'bicycling_duration']

        od_matrix_avg = calculate_avg_time_distance(hexs, 
                                                    od_matrix,
                                                    population='cant_pers')

        indicators_vars = indicators_vars + ['distance_osm_walk_downtown', 
                                             'transit_duration_downtown', 
                                             'driving_duration_in_traffic_downtown',
                                             'walking_duration_downtown', 
                                             'bicycling_duration_downtown']

        print_time_distance(hexs, 
                            od_matrix_avg, 
                            indicators_vars=indicators_vars, 
                            current_path=current_path, 
                            city=city, 
                            prs=prs)

    try:
        print('')
        prs.save(file_pptx)
        print(file_pptx)
        print('')
    except:
        print('No se pudo guardar el archivo', file_pptx)


    if len(od_establecimientos) > 0:
        print('Establecimientos')
        prs = Presentation()
        prs.slide_height = Inches(13.5)
        prs.slide_width = Inches(24)
        file_pptx = current_path / 'Resultados_pptx' / f'{city}_Establecimientos.pptx'


        print_time_distance(hexs, 
                            od_establecimientos, 
                            indicators_vars=['distance', 'duration'],
                            colors_dict={'distance':'Reds', 'duration':'Blues'},
                            equipment_type = equipment_type,
                            current_path=current_path, 
                            city=city,
                            prs=prs)
    
    if len(hexs_green_space) > 0:
        print_time_distance(hexs, 
                    hexs_green_space, 
                    indicators_vars=['green_area_ha_in2000m'],
                    colors_dict={'green_area_ha_in2000m':'Greens'},                    
                    current_path=current_path, 
                    city=city,
                    prs=prs)

        try:
            print('')
            prs.save(file_pptx)
            print(file_pptx)
        except:
            print('No se pudo guardar el archivo', file_pptx)










